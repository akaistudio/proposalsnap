"""
ProposalSnap - AI Presentation Maker
=====================================
Upload your logo, describe your proposal, get a professional PPTX in seconds.
"""

import os, json, uuid, subprocess, colorsys, hashlib, secrets
import bcrypt, smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from datetime import datetime, timedelta
from pathlib import Path
from io import BytesIO
from functools import wraps

import anthropic
import psycopg2
import psycopg2.extras
import requests as http_requests
from PIL import Image
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt
from flask import Flask, request, jsonify, send_file, render_template, render_template_string, redirect, session, flash

app = Flask(__name__)
app.secret_key = os.environ.get('SECRET_KEY', 'proposalsnap-prod-2026')
app.config['PERMANENT_SESSION_LIFETIME'] = timedelta(days=7)
app.config['SESSION_COOKIE_HTTPONLY'] = True
app.config['SESSION_COOKIE_SAMESITE'] = 'Lax'
app.config['MAX_CONTENT_LENGTH'] = 20 * 1024 * 1024

UPLOAD_DIR = Path(__file__).parent / "uploads"
OUTPUT_DIR = Path(__file__).parent / "outputs"
UPLOAD_DIR.mkdir(exist_ok=True)
OUTPUT_DIR.mkdir(exist_ok=True)

# ── Database (lightweight — just users + usage tracking) ────────
def get_db():
    db_url = os.environ.get('DATABASE_URL', '')
    if not db_url:
        return None
    if db_url.startswith('postgres://'):
        db_url = db_url.replace('postgres://', 'postgresql://', 1)
    conn = psycopg2.connect(db_url)
    conn.autocommit = True
    return conn

def init_db():
    try:
        conn = get_db()
        if not conn: return
        cur = conn.cursor()
        cur.execute('''CREATE TABLE IF NOT EXISTS users (
            id SERIAL PRIMARY KEY,
            email TEXT UNIQUE NOT NULL,
            password_hash TEXT NOT NULL,
            company_name TEXT DEFAULT '',
            currency TEXT DEFAULT 'USD',
            is_superadmin BOOLEAN DEFAULT FALSE,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )''')
        cur.execute('''CREATE TABLE IF NOT EXISTS usage_log (
            id SERIAL PRIMARY KEY,
            user_id INTEGER REFERENCES users(id),
            action TEXT DEFAULT 'generate',
            title TEXT DEFAULT '',
            slides INTEGER DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )''')
        # Migrations
        migrations = [
            "ALTER TABLE users ADD COLUMN IF NOT EXISTS is_superadmin BOOLEAN DEFAULT FALSE",
            "UPDATE users SET is_superadmin = TRUE WHERE id = (SELECT MIN(id) FROM users)",
            "ALTER TABLE users ADD COLUMN IF NOT EXISTS brand_logo TEXT DEFAULT ''",
            "ALTER TABLE users ADD COLUMN IF NOT EXISTS brand_colors TEXT DEFAULT ''",
            "ALTER TABLE users ADD COLUMN IF NOT EXISTS brand_font TEXT DEFAULT 'aptos'",
            """CREATE TABLE IF NOT EXISTS otp_codes (
                id SERIAL PRIMARY KEY, email TEXT NOT NULL, code TEXT NOT NULL,
                purpose TEXT DEFAULT 'login', attempts INTEGER DEFAULT 0,
                used BOOLEAN DEFAULT FALSE, created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                expires_at TIMESTAMP NOT NULL)""",
            "CREATE INDEX IF NOT EXISTS idx_otp_email ON otp_codes(email, purpose, used)",
        ]
        for m in migrations:
            try: cur.execute(m)
            except: pass
        conn.close()
        print("✅ Database ready")
    except Exception as e:
        print(f"⚠️ DB not available yet: {e}")

init_db()

# ── Auth helpers ────────────────────────────────────────────────
def hash_pw(pw):
    return bcrypt.hashpw(pw.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')

def check_pw(pw, hashed):
    try:
        return bcrypt.checkpw(pw.encode('utf-8'), hashed.encode('utf-8'))
    except (ValueError, AttributeError):
        if hashlib.sha256(pw.encode()).hexdigest() == hashed:
            return True
        return False

def generate_otp():
    return f"{secrets.randbelow(900000) + 100000}"

def send_otp_email(email, code, purpose='login'):
    resend_key = os.environ.get('RESEND_API_KEY', '')
    from_email = os.environ.get('SMTP_FROM', 'onboarding@resend.dev')
    purpose_text = 'login' if purpose == 'login' else 'verification'
    subject = f"Your ProposalSnap {purpose_text} code: {code}"
    html = f"""<div style="font-family:sans-serif;max-width:400px;margin:0 auto;padding:24px">
        <h2 style="color:#f59e0b">ProposalSnap</h2>
        <p style="color:#666;font-size:14px">Your {purpose_text} code is:</p>
        <div style="font-size:36px;font-weight:800;letter-spacing:8px;color:#1a1a2e;text-align:center;
                    padding:20px;background:#f0f4ff;border-radius:12px;margin:16px 0">{code}</div>
        <p style="color:#999;font-size:12px">This code expires in 5 minutes. Do not share it.</p>
        <p style="color:#999;font-size:11px;margin-top:20px">Part of <a href="https://snapsuite.up.railway.app" style="color:#f59e0b">SnapSuite</a></p>
    </div>"""
    if not resend_key:
        print(f"⚠️ RESEND_API_KEY not set. OTP for {email}: {code}")
        return True
    import requests as http_req
    try:
        r = http_req.post('https://api.resend.com/emails', json={
            'from': from_email, 'to': [email], 'subject': subject, 'html': html
        }, headers={'Authorization': f'Bearer {resend_key}'}, timeout=10)
        if r.status_code == 200:
            print(f"✅ OTP sent to {email}")
            return True
        else:
            print(f"❌ Resend error {r.status_code}: {r.text}")
            print(f"💡 OTP for {email}: {code}")
            return True
    except Exception as e:
        print(f"❌ Email failed: {e}")
        print(f"💡 OTP for {email}: {code}")
        return True


def register_with_hub(company_name, email, currency):
    hub = os.environ.get('FINANCESNAP_URL', 'https://snapsuite.up.railway.app')
    try:
        http_requests.post(f'{hub}/api/register-company', json={
            'app_name': 'ProposalSnap', 'company_name': company_name,
            'email': email, 'currency': currency,
            'app_url': 'https://proposalsnap.up.railway.app'
        }, timeout=5)
    except: pass

def login_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if 'user_id' not in session:
            return redirect('/login')
        return f(*args, **kwargs)
    return decorated

def log_usage(title='', slides=0):
    """Log a generation event for the current user"""
    try:
        uid = session.get('user_id')
        if not uid: return
        conn = get_db()
        if not conn: return
        cur = conn.cursor()
        cur.execute('INSERT INTO usage_log (user_id, action, title, slides) VALUES (%s,%s,%s,%s)',
                    (uid, 'generate', title[:200], slides))
        conn.close()
    except: pass

MODEL = "claude-sonnet-4-5-20250929"
client = anthropic.Anthropic(api_key=os.environ.get('ANTHROPIC_API_KEY', ''))

# ── Color Extraction ──────────────────────────────────────────
def extract_colors_from_logo(logo_path):
    """Extract dominant colors from logo and generate complementary palette"""
    try:
        img = Image.open(logo_path).convert("RGB")
        img = img.resize((100, 100))
        pixels = list(img.getdata())
        
        # Filter out near-white and near-black pixels
        filtered = [(r, g, b) for r, g, b in pixels 
                     if not (r > 240 and g > 240 and b > 240) 
                     and not (r < 15 and g < 15 and b < 15)]
        
        if not filtered:
            filtered = pixels
        
        # Simple dominant color: average of filtered pixels
        avg_r = int(sum(p[0] for p in filtered) / len(filtered))
        avg_g = int(sum(p[1] for p in filtered) / len(filtered))
        avg_b = int(sum(p[2] for p in filtered) / len(filtered))
        
        # Get most vivid pixel (highest saturation)
        best_sat = 0
        vivid = (avg_r, avg_g, avg_b)
        for r, g, b in filtered:
            h, s, v = colorsys.rgb_to_hsv(r/255, g/255, b/255)
            if s > best_sat and v > 0.2:
                best_sat = s
                vivid = (r, g, b)
        
        primary_hex = f"{vivid[0]:02x}{vivid[1]:02x}{vivid[2]:02x}"
        
        # Generate complementary colors
        h, s, v = colorsys.rgb_to_hsv(vivid[0]/255, vivid[1]/255, vivid[2]/255)
        
        # Secondary: lighter version
        sr, sg, sb = colorsys.hsv_to_rgb(h, max(0.1, s * 0.3), min(1.0, v * 1.4 + 0.3))
        secondary_hex = f"{int(sr*255):02x}{int(sg*255):02x}{int(sb*255):02x}"
        
        # Accent: shifted hue, high saturation
        ah, as_, av = (h + 0.08) % 1.0, min(1.0, s * 1.2 + 0.1), min(1.0, v * 1.1 + 0.1)
        ar, ag, ab = colorsys.hsv_to_rgb(ah, as_, av)
        accent_hex = f"{int(ar*255):02x}{int(ag*255):02x}{int(ab*255):02x}"
        
        # Dark: very dark version of primary
        dr, dg, db = colorsys.hsv_to_rgb(h, min(0.6, s * 0.8), 0.12)
        dark_hex = f"{int(dr*255):02x}{int(dg*255):02x}{int(db*255):02x}"
        
        return {
            "primary": primary_hex,
            "secondary": secondary_hex,
            "accent": accent_hex,
            "dark": dark_hex,
            "light": "F8F9FA",
            "textDark": "1A1A2E",
            "textLight": "FFFFFF",
            "textMuted": "6B7280"
        }
    except Exception as e:
        print(f"Color extraction error: {e}")
        return default_colors()

def default_colors():
    return {
        "primary": "1E2761", "secondary": "CADCFC", "accent": "4A90D9",
        "dark": "0F1629", "light": "F8F9FA", "textDark": "1A1A2E",
        "textLight": "FFFFFF", "textMuted": "6B7280"
    }

# ── Claude API ────────────────────────────────────────────────
def generate_slide_content(client_name, company_name, pres_type, tone, key_points, num_slides=12):
    """Use Claude to generate structured slide content"""
    prompt = f"""Generate a professional {pres_type} presentation structure with STRONG visual variety.

Client: {client_name}
Presenting Company: {company_name}
Tone: {tone}
Number of Slides: {num_slides}

Key Points / Requirements:
{key_points}

Return ONLY a valid JSON array of slide objects. Each slide MUST have these fields:
- "layout": one of the layouts below
- "title": slide title
- Additional fields based on layout:

AVAILABLE LAYOUTS:

"title": Opening slide. Fields: subtitle (string)
"agenda": Overview of what's covered. Fields: bullets (array of 5-7 strings)
"content": Standard text slide. Fields: bullets (array of 3-5 strings) OR body (paragraph), optional subtitle
"two_column": Side-by-side comparison. Fields: left_title, left_bullets (array), right_title, right_bullets (array)
"stats": Big number metrics (2-4 cards). Fields: stats (array of {{value, label, description}})
"timeline": Process/timeline steps. Fields: steps (array of {{phase, description, duration}})
"pricing": Investment/tier cards. Fields: tiers (array of {{name, price, features[], highlight bool}})
"team": Team member cards. Fields: members (array of {{name, role, bio}})
"icon_grid": 4-6 feature/service cards with icons. Fields: items (array of {{icon, heading, description}}). icon should be a single emoji.
"comparison": Before vs After or comparison table. Fields: left_label, right_label, rows (array of {{feature, left_value, right_value}})
"quote": Testimonial or key statement. Fields: quote (string), attribution (string), role (string)
"metric_bar": Horizontal progress bars. Fields: metrics (array of {{label, value, max_value, description}}). value and max_value are numbers (e.g. value:85, max_value:100)
"process_flow": Numbered process with arrows. Fields: steps (array of {{number, title, description}})
"checklist": Visual checklist/deliverables. Fields: items (array of strings), subtitle (string)
"big_statement": One powerful sentence in large text. Fields: statement (string), supporting_text (string)
"closing": Thank you slide. Fields: subtitle, contact (string)

RULES:
1. If generating 4+ slides: First slide MUST be "title", second MUST be "agenda", last MUST be "closing"
2. If generating 1-3 slides: Use the most impactful layouts. For 1 slide use "big_statement" or "content". For 2-3 slides, start with "title" and use visual layouts for the rest. No agenda or closing needed.
3. CRITICAL: Use as many different layout types as possible. Do NOT repeat the same layout more than twice.
4. For 6+ slides: ALWAYS include at least one "stats", one "timeline" or "process_flow", and one of "icon_grid"/"comparison"/"metric_bar"
5. Prefer visual layouts (stats, icon_grid, comparison, metric_bar, process_flow, checklist, quote) over plain "content"
6. Use "content" for maximum 2 slides. Use visual layouts for the rest.
7. Stats should have realistic, specific numbers
8. Bullets should be concise (10-20 words each)
9. Make content specific to the client and key points
10. Return ONLY the JSON array, no other text

Generate exactly {num_slides} slides."""

    response = client.messages.create(
        model=MODEL, max_tokens=4000,
        messages=[{"role": "user", "content": prompt}]
    )
    text = response.content[0].text.strip()
    if text.startswith("```"):
        text = text.split("\n", 1)[1] if "\n" in text else text[3:]
        text = text.rsplit("```", 1)[0]
    return json.loads(text)

# ── Extract slides from uploaded PPTX ──────────────────────────
def extract_slides_from_pptx(filepath):
    """Extract text content from each slide of an uploaded PPTX"""
    prs = PptxPresentation(filepath)
    slides = []
    for i, slide in enumerate(prs.slides):
        slide_data = {"slide_number": i + 1, "texts": []}
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_data["texts"].append(text)
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                slide_data["table"] = rows
        if slide_data["texts"] or slide_data.get("table"):
            slides.append(slide_data)
    return slides

def polish_slide_content(original_slides, instructions, num_slides=None, tone="Corporate"):
    """Use Claude to polish/improve existing slide content"""
    slides_text = json.dumps(original_slides, indent=2)
    if num_slides is None:
        num_slides = len(original_slides)

    prompt = f"""You are a presentation expert. I have an existing presentation with the following slide content:

{slides_text}

USER'S INSTRUCTIONS FOR POLISHING:
{instructions}

Take the existing content and POLISH it according to the user's instructions.

Return ONLY a valid JSON array of slide objects. Each slide MUST have these fields:
- "layout": one of the layouts below
- "title": slide title
- Additional fields based on layout:

AVAILABLE LAYOUTS:

"title": Opening slide. Fields: subtitle (string)
"agenda": Overview of what's covered. Fields: bullets (array of 5-7 strings)
"content": Standard text slide. Fields: bullets (array of 3-5 strings) OR body (paragraph), optional subtitle
"two_column": Side-by-side comparison. Fields: left_title, left_bullets (array), right_title, right_bullets (array)
"stats": Big number metrics (2-4 cards). Fields: stats (array of {{value, label, description}})
"timeline": Process/timeline steps. Fields: steps (array of {{phase, description, duration}})
"pricing": Investment/tier cards. Fields: tiers (array of {{name, price, features[], highlight bool}})
"team": Team member cards. Fields: members (array of {{name, role, bio}})
"icon_grid": 4-6 feature/service cards with icons. Fields: items (array of {{icon, heading, description}}). icon should be a single emoji.
"comparison": Before vs After or comparison table. Fields: left_label, right_label, rows (array of {{feature, left_value, right_value}})
"quote": Testimonial or key statement. Fields: quote (string), attribution (string), role (string)
"metric_bar": Horizontal progress bars. Fields: metrics (array of {{label, value, max_value, description}}). value and max_value are numbers.
"process_flow": Numbered process with arrows. Fields: steps (array of {{number, title, description}})
"checklist": Visual checklist/deliverables. Fields: items (array of strings), subtitle (string)
"big_statement": One powerful sentence in large text. Fields: statement (string), supporting_text (string)
"closing": Thank you slide. Fields: subtitle, contact (string)

RULES:
1. PRESERVE the original content and meaning — polish, don't rewrite from scratch
2. Apply the user's specific instructions (tone changes, additions, restructuring, etc.)
3. Use visual layouts (stats, icon_grid, comparison, timeline) where the content fits
4. For 4+ slides: First slide should be "title", last should be "closing"
5. For 1-3 slides: Use the most impactful layouts, no need for title/closing wrapper
6. Use as many different layout types as possible for visual variety
7. Keep bullets concise (10-20 words each)
8. Generate exactly {num_slides} slides
9. Tone: {tone}
10. Return ONLY the JSON array, no other text"""

    response = client.messages.create(
        model=MODEL, max_tokens=4000,
        messages=[{"role": "user", "content": prompt}]
    )
    text = response.content[0].text.strip()
    if text.startswith("```"):
        text = text.split("\n", 1)[1] if "\n" in text else text[3:]
        text = text.rsplit("```", 1)[0]
    return json.loads(text)

# ── Brand Management ───────────────────────────────────────────
def get_user_brand():
    """Load saved brand settings for current user"""
    try:
        uid = session.get('user_id')
        if not uid: return None
        conn = get_db()
        if not conn: return None
        cur = conn.cursor(cursor_factory=psycopg2.extras.RealDictCursor)
        cur.execute('SELECT brand_logo, brand_colors, brand_font FROM users WHERE id=%s', (uid,))
        user = cur.fetchone()
        conn.close()
        if not user: return None
        brand = {}
        if user.get('brand_colors'):
            try: brand['colors'] = json.loads(user['brand_colors'])
            except: pass
        if user.get('brand_logo') and Path(user['brand_logo']).exists():
            brand['logo_path'] = user['brand_logo']
        if user.get('brand_font'):
            brand['font'] = user['brand_font']
        return brand if brand else None
    except: return None

def apply_brand(colors, logo_path, font_style):
    """Override with saved brand if user hasn't provided custom ones"""
    brand = get_user_brand()
    if not brand: return colors, logo_path, font_style
    if not logo_path and brand.get('logo_path'):
        logo_path = brand['logo_path']
    if colors == default_colors() and brand.get('colors'):
        colors = brand['colors']
    if font_style == 'aptos' and brand.get('font'):
        font_style = brand['font']
    return colors, logo_path, font_style

# ── Audience Versioning ───────────────────────────────────────
def generate_audience_version(original_slides, audience_type, num_slides):
    """Generate a version of the deck targeted at a specific audience"""
    slides_text = json.dumps(original_slides, indent=2)

    audiences = {
        "executive": {
            "desc": "C-suite executives with limited time",
            "rules": "Maximum 6-8 slides. Lead with conclusion/recommendation. Use stats and big statements. Cut all technical details. Every slide must be answerable in 30 seconds.",
            "default_slides": 6
        },
        "detailed": {
            "desc": "Team members and managers who need full context",
            "rules": "12-16 slides. Include all details, process flows, timelines, and supporting data. Use a mix of visual layouts. Keep the narrative comprehensive.",
            "default_slides": 14
        },
        "investor": {
            "desc": "Investors and board members evaluating opportunity",
            "rules": "8-10 slides. Follow: Problem → Solution → Market → Traction → Model → Team → Ask. Make numbers and metrics prominent. Include comparison/competitive landscape. End with clear ask/next steps.",
            "default_slides": 10
        }
    }

    cfg = audiences.get(audience_type, audiences["detailed"])
    slide_count = num_slides or cfg["default_slides"]

    prompt = f"""You are a presentation expert creating a version of a deck for: {cfg["desc"]}

ORIGINAL CONTENT:
{slides_text}

AUDIENCE RULES:
{cfg["rules"]}

Restructure and adapt the content for this audience. Generate exactly {slide_count} slides.

Return ONLY a valid JSON array of slide objects. Each slide MUST have:
- "layout": one of: title, agenda, content, two_column, stats, timeline, pricing, team, icon_grid, comparison, quote, metric_bar, process_flow, checklist, big_statement, closing
- "title": slide title
- Layout-specific fields (bullets, stats, steps, etc.)

RULES:
1. First slide = "title", last = "closing"
2. Use at LEAST 5 different layout types
3. Prefer visual layouts (stats, icon_grid, comparison, timeline, metric_bar) over plain "content"
4. Keep bullets concise (10-20 words)
5. Return ONLY the JSON array"""

    response = client.messages.create(
        model=MODEL, max_tokens=4000,
        messages=[{"role": "user", "content": prompt}]
    )
    text = response.content[0].text.strip()
    if text.startswith("```"):
        text = text.split("\n", 1)[1] if "\n" in text else text[3:]
        text = text.rsplit("```", 1)[0]
    return json.loads(text)

# ── Style Transfer ────────────────────────────────────────────
def extract_style_from_pptx(filepath):
    """Extract style patterns (structure, tone, layout types) from a reference deck"""
    prs = PptxPresentation(filepath)
    style_info = {
        "total_slides": len(prs.slides),
        "slide_structures": [],
        "color_samples": [],
        "font_samples": [],
        "tone_samples": []
    }

    for i, slide in enumerate(prs.slides):
        structure = {"slide_num": i + 1, "shape_count": len(slide.shapes), "has_table": False,
                     "has_image": False, "text_blocks": 0, "texts": []}
        for shape in slide.shapes:
            if shape.has_text_frame:
                structure["text_blocks"] += 1
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        structure["texts"].append(text[:100])
                    for run in para.runs:
                        if run.font.color and run.font.color.rgb:
                            style_info["color_samples"].append(str(run.font.color.rgb))
                        if run.font.name:
                            style_info["font_samples"].append(run.font.name)
            if shape.has_table:
                structure["has_table"] = True
            if shape.shape_type == 13:  # picture
                structure["has_image"] = True
        style_info["slide_structures"].append(structure)
        if structure["texts"]:
            style_info["tone_samples"].extend(structure["texts"][:2])

    return style_info

def style_transfer_content(content_slides, style_info, instructions=""):
    """Apply extracted style patterns to content"""
    content_text = json.dumps(content_slides, indent=2)
    style_text = json.dumps(style_info, indent=2)

    prompt = f"""You are a presentation design expert. Apply the STYLE from a reference deck to the CONTENT of another deck.

REFERENCE DECK STYLE:
{style_text}

CONTENT TO RESTYLE:
{content_text}

{"ADDITIONAL INSTRUCTIONS: " + instructions if instructions else ""}

Analyze the reference deck's patterns:
- How many slides and what structure (heavy on data? visual? minimal text?)
- Tone of the text (formal? conversational? punchy?)
- Use of tables, images, comparisons
- Slide flow and organization

Then restructure the content to MATCH that style while preserving the content's meaning.

Return ONLY a valid JSON array of slide objects. Each slide MUST have:
- "layout": one of: title, agenda, content, two_column, stats, timeline, pricing, team, icon_grid, comparison, quote, metric_bar, process_flow, checklist, big_statement, closing
- "title": slide title
- Layout-specific fields

RULES:
1. Match the reference deck's structure and rhythm
2. If reference is concise (6-8 slides), output the same count. If detailed (15+), expand.
3. If reference uses lots of data/stats, convert content into data-heavy layouts
4. If reference is text-minimal, keep bullets very short
5. First slide = "title", last = "closing"
6. Use at LEAST 5 different layout types
7. Return ONLY the JSON array"""

    response = client.messages.create(
        model=MODEL, max_tokens=4000,
        messages=[{"role": "user", "content": prompt}]
    )
    text = response.content[0].text.strip()
    if text.startswith("```"):
        text = text.split("\n", 1)[1] if "\n" in text else text[3:]
        text = text.rsplit("```", 1)[0]
    return json.loads(text)

# ── PPTX Generation ───────────────────────────────────────────
def create_pptx(slides, colors, client_name, company_name, pres_type, tone, logo_path=None, font_style='aptos'):
    """Generate PPTX using Node.js pptxgenjs"""
    output_id = str(uuid.uuid4())[:8]
    output_path = str(OUTPUT_DIR / f"proposal_{output_id}.pptx")
    
    input_data = {
        "outputPath": output_path,
        "clientName": client_name,
        "companyName": company_name,
        "presentationType": pres_type,
        "tone": tone,
        "slides": slides,
        "colors": colors,
        "logoPath": str(logo_path) if logo_path else None,
        "fontStyle": font_style
    }
    
    script_path = Path(__file__).parent / "generate_pptx.js"
    
    # Find node binary
    import shutil
    node_bin = shutil.which("node")
    if not node_bin:
        # Try common paths
        for p in ["/usr/bin/node", "/usr/local/bin/node", "/nix/store/*/bin/node"]:
            import glob
            matches = glob.glob(p)
            if matches:
                node_bin = matches[0]
                break
    if not node_bin:
        node_bin = "node"
    
    result = subprocess.run(
        [node_bin, str(script_path)],
        input=json.dumps(input_data),
        capture_output=True, text=True, timeout=30
    )
    
    if result.returncode != 0:
        raise Exception(f"PPTX generation failed: {result.stderr}")
    
    return output_path

# ── API Routes ────────────────────────────────────────────────
@app.route('/api/generate', methods=['POST'])
def generate():
    try:
        client_name = request.form.get('client_name', '').strip()
        company_name = request.form.get('company_name', '').strip()
        pres_type = request.form.get('presentation_type', 'Corporate Proposal')
        tone = request.form.get('tone', 'Corporate')
        font_style = request.form.get('font_style', 'aptos')
        key_points = request.form.get('key_points', '').strip()
        num_slides = int(request.form.get('num_slides', 12))
        num_slides = max(1, min(16, num_slides))
        
        if not client_name: return jsonify({"error": "Client name is required"}), 400
        if not key_points: return jsonify({"error": "Key points are required"}), 400
        
        # Handle logo
        logo_path = None
        colors = default_colors()
        if 'logo' in request.files and request.files['logo'].filename:
            logo_file = request.files['logo']
            logo_ext = Path(logo_file.filename).suffix.lower()
            if logo_ext in ('.png', '.jpg', '.jpeg', '.svg', '.webp'):
                logo_id = str(uuid.uuid4())[:8]
                logo_path = UPLOAD_DIR / f"logo_{logo_id}{logo_ext}"
                logo_file.save(str(logo_path))
                if logo_ext != '.svg':
                    colors = extract_colors_from_logo(str(logo_path))
        
        # Generate content with Claude
        slides = generate_slide_content(client_name, company_name, pres_type, tone, key_points, num_slides)
        
        # Apply saved brand if no custom logo/colors provided
        colors, logo_path, font_style = apply_brand(colors, logo_path, font_style)
        
        # Generate PPTX
        output_path = create_pptx(slides, colors, client_name, company_name, pres_type, tone, logo_path, font_style)
        
        filename = f"{client_name.replace(' ', '_')}_{pres_type.replace(' ', '_')}.pptx"
        log_usage(title=f"{client_name} — {pres_type}", slides=len(slides))
        return jsonify({
            "success": True,
            "download_url": f"/api/download/{Path(output_path).name}",
            "filename": filename,
            "slides_count": len(slides),
            "colors": colors
        })
    except json.JSONDecodeError:
        return jsonify({"error": "Failed to generate slide content. Please try again."}), 500
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/polish', methods=['POST'])
def polish():
    try:
        instructions = request.form.get('instructions', '').strip()
        tone = request.form.get('tone', 'Corporate')
        font_style = request.form.get('font_style', 'aptos')
        num_slides = request.form.get('num_slides', '')

        if not instructions:
            return jsonify({"error": "Please provide polishing instructions"}), 400

        if 'pptx_file' not in request.files or not request.files['pptx_file'].filename:
            return jsonify({"error": "Please upload a PPTX file"}), 400

        pptx_file = request.files['pptx_file']
        if not pptx_file.filename.lower().endswith('.pptx'):
            return jsonify({"error": "Only .pptx files are supported"}), 400

        # Save uploaded file
        upload_id = str(uuid.uuid4())[:8]
        upload_path = UPLOAD_DIR / f"upload_{upload_id}.pptx"
        pptx_file.save(str(upload_path))

        # Extract content from uploaded slides
        original_slides = extract_slides_from_pptx(str(upload_path))
        if not original_slides:
            return jsonify({"error": "Could not extract any content from the uploaded file"}), 400

        # Determine slide count
        slide_count = int(num_slides) if num_slides else len(original_slides)
        slide_count = max(1, min(20, slide_count))

        # Polish with AI
        polished_slides = polish_slide_content(original_slides, instructions, slide_count, tone)

        # Handle logo if provided
        logo_path = None
        colors = default_colors()
        if 'logo' in request.files and request.files['logo'].filename:
            logo_file = request.files['logo']
            logo_ext = Path(logo_file.filename).suffix.lower()
            if logo_ext in ('.png', '.jpg', '.jpeg', '.webp', '.svg'):
                logo_id = str(uuid.uuid4())[:8]
                logo_path = UPLOAD_DIR / f"logo_{logo_id}{logo_ext}"
                logo_file.save(str(logo_path))
                if logo_ext != '.svg':
                    colors = extract_colors_from_logo(str(logo_path))

        # Get client/company from first slide title or form
        client_name = request.form.get('client_name', '').strip()
        if not client_name and polished_slides:
            client_name = polished_slides[0].get('title', 'Polished Deck')
        company_name = request.form.get('company_name', '').strip()

        # Apply saved brand if no custom logo/colors provided
        colors, logo_path, font_style = apply_brand(colors, logo_path, font_style)

        # Generate polished PPTX
        output_path = create_pptx(polished_slides, colors, client_name, company_name,
                                  'Polished Presentation', tone, logo_path, font_style)

        filename = f"Polished_{client_name.replace(' ', '_')}.pptx" if client_name else "Polished_Presentation.pptx"
        log_usage(title=f"Polish: {client_name}", slides=len(polished_slides))

        return jsonify({
            "success": True,
            "download_url": f"/api/download/{Path(output_path).name}",
            "filename": filename,
            "slides_count": len(polished_slides),
            "original_slides": len(original_slides),
            "colors": colors
        })
    except json.JSONDecodeError:
        return jsonify({"error": "AI failed to generate valid slide structure. Please try again."}), 500
    except Exception as e:
        return jsonify({"error": str(e)}), 500

        return jsonify({"error": str(e)}), 500

# ── Brand Settings ────────────────────────────────────────────
@app.route('/api/brand', methods=['GET', 'POST'])
def brand_settings():
    uid = session.get('user_id')
    if not uid: return jsonify({"error": "Not logged in"}), 401

    conn = get_db()
    if not conn: return jsonify({"error": "Database not configured"}), 500
    cur = conn.cursor(cursor_factory=psycopg2.extras.RealDictCursor)

    if request.method == 'GET':
        cur.execute('SELECT brand_logo, brand_colors, brand_font, company_name FROM users WHERE id=%s', (uid,))
        user = cur.fetchone()
        conn.close()
        colors = {}
        try: colors = json.loads(user.get('brand_colors') or '{}')
        except: pass
        return jsonify({"colors": colors, "font": user.get('brand_font', 'aptos'),
                       "has_logo": bool(user.get('brand_logo')), "company": user.get('company_name', '')})

    # POST — save brand
    font = request.form.get('font', 'aptos')
    cur.execute('UPDATE users SET brand_font=%s WHERE id=%s', (font, uid))

    if 'logo' in request.files and request.files['logo'].filename:
        logo_file = request.files['logo']
        logo_ext = Path(logo_file.filename).suffix.lower()
        if logo_ext in ('.png', '.jpg', '.jpeg', '.webp', '.svg'):
            logo_id = str(uuid.uuid4())[:8]
            logo_path = UPLOAD_DIR / f"brand_{uid}_{logo_id}{logo_ext}"
            logo_file.save(str(logo_path))
            colors = default_colors()
            if logo_ext != '.svg':
                colors = extract_colors_from_logo(str(logo_path))
            cur.execute('UPDATE users SET brand_logo=%s, brand_colors=%s WHERE id=%s',
                       (str(logo_path), json.dumps(colors), uid))

    conn.close()
    return jsonify({"success": True})

# ── Audience Versioning ───────────────────────────────────────
@app.route('/api/version', methods=['POST'])
def version():
    try:
        if 'pptx_file' not in request.files or not request.files['pptx_file'].filename:
            return jsonify({"error": "Please upload a PPTX file"}), 400

        pptx_file = request.files['pptx_file']
        upload_id = str(uuid.uuid4())[:8]
        upload_path = UPLOAD_DIR / f"upload_{upload_id}.pptx"
        pptx_file.save(str(upload_path))

        original_slides = extract_slides_from_pptx(str(upload_path))
        if not original_slides:
            return jsonify({"error": "Could not extract content from file"}), 400

        audiences = request.form.getlist('audiences') or ['executive', 'detailed', 'investor']
        font_style = request.form.get('font_style', 'aptos')
        client_name = request.form.get('client_name', '').strip() or 'Presentation'
        company_name = request.form.get('company_name', '').strip()

        # Handle logo / brand
        logo_path = None
        colors = default_colors()
        if 'logo' in request.files and request.files['logo'].filename:
            logo_file = request.files['logo']
            logo_ext = Path(logo_file.filename).suffix.lower()
            if logo_ext in ('.png', '.jpg', '.jpeg', '.webp', '.svg'):
                lid = str(uuid.uuid4())[:8]
                logo_path = UPLOAD_DIR / f"logo_{lid}{logo_ext}"
                logo_file.save(str(logo_path))
                if logo_ext != '.svg':
                    colors = extract_colors_from_logo(str(logo_path))
        colors, logo_path, font_style = apply_brand(colors, logo_path, font_style)

        results = []
        for audience in audiences:
            try:
                slides = generate_audience_version(original_slides, audience, num_slides=None)
                output_path = create_pptx(slides, colors, client_name, company_name,
                                         f'{audience.title()} Version', 'Corporate',
                                         logo_path, font_style)
                fname = f"{client_name.replace(' ', '_')}_{audience}.pptx"
                log_usage(title=f"Version: {audience} — {client_name}", slides=len(slides))
                results.append({
                    "audience": audience,
                    "download_url": f"/api/download/{Path(output_path).name}",
                    "filename": fname,
                    "slides_count": len(slides)
                })
            except Exception as e:
                results.append({"audience": audience, "error": str(e)})

        return jsonify({"success": True, "versions": results, "original_slides": len(original_slides)})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# ── Style Transfer ────────────────────────────────────────────
@app.route('/api/style-transfer', methods=['POST'])
def style_transfer():
    try:
        if 'content_file' not in request.files or not request.files['content_file'].filename:
            return jsonify({"error": "Please upload your content PPTX"}), 400
        if 'reference_file' not in request.files or not request.files['reference_file'].filename:
            return jsonify({"error": "Please upload a reference/style PPTX"}), 400

        uid = str(uuid.uuid4())[:8]

        content_file = request.files['content_file']
        content_path = UPLOAD_DIR / f"content_{uid}.pptx"
        content_file.save(str(content_path))

        ref_file = request.files['reference_file']
        ref_path = UPLOAD_DIR / f"ref_{uid}.pptx"
        ref_file.save(str(ref_path))

        instructions = request.form.get('instructions', '').strip()
        font_style = request.form.get('font_style', 'aptos')
        client_name = request.form.get('client_name', '').strip() or 'Styled Deck'
        company_name = request.form.get('company_name', '').strip()

        content_slides = extract_slides_from_pptx(str(content_path))
        if not content_slides:
            return jsonify({"error": "Could not extract content from your deck"}), 400

        style_info = extract_style_from_pptx(str(ref_path))

        styled_slides = style_transfer_content(content_slides, style_info, instructions)

        # Handle logo / brand
        logo_path = None
        colors = default_colors()
        if 'logo' in request.files and request.files['logo'].filename:
            logo_file = request.files['logo']
            logo_ext = Path(logo_file.filename).suffix.lower()
            if logo_ext in ('.png', '.jpg', '.jpeg', '.webp', '.svg'):
                lid = str(uuid.uuid4())[:8]
                logo_path = UPLOAD_DIR / f"logo_{lid}{logo_ext}"
                logo_file.save(str(logo_path))
                if logo_ext != '.svg':
                    colors = extract_colors_from_logo(str(logo_path))
        colors, logo_path, font_style = apply_brand(colors, logo_path, font_style)

        output_path = create_pptx(styled_slides, colors, client_name, company_name,
                                  'Styled Presentation', 'Corporate', logo_path, font_style)

        fname = f"Styled_{client_name.replace(' ', '_')}.pptx"
        log_usage(title=f"Style Transfer: {client_name}", slides=len(styled_slides))

        return jsonify({
            "success": True,
            "download_url": f"/api/download/{Path(output_path).name}",
            "filename": fname,
            "slides_count": len(styled_slides),
            "reference_slides": style_info["total_slides"],
            "content_slides": len(content_slides)
        })
    except json.JSONDecodeError:
        return jsonify({"error": "AI failed to generate valid structure. Please try again."}), 500
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# ── Merge Decks ───────────────────────────────────────────────
@app.route('/api/merge', methods=['POST'])
def merge_decks():
    try:
        files = request.files.getlist('pptx_files')
        if len(files) < 2:
            return jsonify({"error": "Please upload at least 2 PPTX files to merge"}), 400

        all_slides = []
        file_names = []
        for f in files:
            if not f.filename.lower().endswith('.pptx'): continue
            uid = str(uuid.uuid4())[:8]
            fpath = UPLOAD_DIR / f"merge_{uid}.pptx"
            f.save(str(fpath))
            slides = extract_slides_from_pptx(str(fpath))
            all_slides.extend(slides)
            file_names.append(f.filename)

        if not all_slides:
            return jsonify({"error": "Could not extract content from uploaded files"}), 400

        instructions = request.form.get('instructions', '').strip()
        if not instructions:
            instructions = "Combine all slide content into one cohesive presentation. Remove duplicate content. Create a unified narrative flow with a single title slide and closing slide."

        font_style = request.form.get('font_style', 'aptos')
        client_name = request.form.get('client_name', '').strip() or 'Merged Deck'
        company_name = request.form.get('company_name', '').strip()
        tone = request.form.get('tone', 'Corporate')

        merged_slides = polish_slide_content(all_slides, instructions, num_slides=None, tone=tone)

        logo_path = None
        colors = default_colors()
        if 'logo' in request.files and request.files['logo'].filename:
            logo_file = request.files['logo']
            logo_ext = Path(logo_file.filename).suffix.lower()
            if logo_ext in ('.png', '.jpg', '.jpeg', '.webp', '.svg'):
                lid = str(uuid.uuid4())[:8]
                logo_path = UPLOAD_DIR / f"logo_{lid}{logo_ext}"
                logo_file.save(str(logo_path))
                if logo_ext != '.svg':
                    colors = extract_colors_from_logo(str(logo_path))
        colors, logo_path, font_style = apply_brand(colors, logo_path, font_style)

        output_path = create_pptx(merged_slides, colors, client_name, company_name,
                                  'Merged Presentation', tone, logo_path, font_style)

        fname = f"Merged_{client_name.replace(' ', '_')}.pptx"
        log_usage(title=f"Merge: {' + '.join(file_names)}", slides=len(merged_slides))

        return jsonify({
            "success": True,
            "download_url": f"/api/download/{Path(output_path).name}",
            "filename": fname,
            "slides_count": len(merged_slides),
            "source_files": len(files),
            "source_slides": len(all_slides)
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# ── Split Deck ────────────────────────────────────────────────
@app.route('/api/split', methods=['POST'])
def split_deck():
    try:
        if 'pptx_file' not in request.files or not request.files['pptx_file'].filename:
            return jsonify({"error": "Please upload a PPTX file"}), 400

        pptx_file = request.files['pptx_file']
        uid = str(uuid.uuid4())[:8]
        upload_path = UPLOAD_DIR / f"split_{uid}.pptx"
        pptx_file.save(str(upload_path))

        original_slides = extract_slides_from_pptx(str(upload_path))
        if not original_slides:
            return jsonify({"error": "Could not extract content from file"}), 400

        font_style = request.form.get('font_style', 'aptos')
        client_name = request.form.get('client_name', '').strip() or 'Presentation'
        company_name = request.form.get('company_name', '').strip()

        logo_path = None
        colors = default_colors()
        if 'logo' in request.files and request.files['logo'].filename:
            logo_file = request.files['logo']
            logo_ext = Path(logo_file.filename).suffix.lower()
            if logo_ext in ('.png', '.jpg', '.jpeg', '.webp', '.svg'):
                lid = str(uuid.uuid4())[:8]
                logo_path = UPLOAD_DIR / f"logo_{lid}{logo_ext}"
                logo_file.save(str(logo_path))
                if logo_ext != '.svg':
                    colors = extract_colors_from_logo(str(logo_path))
        colors, logo_path, font_style = apply_brand(colors, logo_path, font_style)

        # Generate executive summary (short) + appendix (detailed)
        results = []

        # Part 1: Executive Summary
        summary_prompt = "Extract only the most critical content. Create a tight 6-slide executive summary: title, key takeaway, 2-3 main points with stats/visuals, recommendation, closing. Be ruthlessly concise."
        summary_slides = polish_slide_content(original_slides, summary_prompt, num_slides=6, tone="Corporate")
        summary_path = create_pptx(summary_slides, colors, client_name, company_name,
                                   'Executive Summary', 'Corporate', logo_path, font_style)
        results.append({
            "part": "Executive Summary",
            "download_url": f"/api/download/{Path(summary_path).name}",
            "filename": f"{client_name.replace(' ', '_')}_Summary.pptx",
            "slides_count": len(summary_slides)
        })

        # Part 2: Full Detail / Appendix
        detail_prompt = "Expand all supporting details into a comprehensive appendix deck. Include all data, processes, timelines, team info, and supporting evidence. This is the deep-dive version."
        detail_slides = polish_slide_content(original_slides, detail_prompt,
                                            num_slides=max(len(original_slides), 12), tone="Corporate")
        detail_path = create_pptx(detail_slides, colors, client_name, company_name,
                                  'Full Detail', 'Corporate', logo_path, font_style)
        results.append({
            "part": "Full Detail + Appendix",
            "download_url": f"/api/download/{Path(detail_path).name}",
            "filename": f"{client_name.replace(' ', '_')}_Detail.pptx",
            "slides_count": len(detail_slides)
        })

        log_usage(title=f"Split: {client_name}", slides=len(summary_slides) + len(detail_slides))

        return jsonify({"success": True, "parts": results, "original_slides": len(original_slides)})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/download/<filename>')
def download(filename):
    filepath = OUTPUT_DIR / filename
    if not filepath.exists():
        return jsonify({"error": "File not found"}), 404
    return send_file(str(filepath), as_attachment=True,
                     download_name=request.args.get('name', filename))

@app.route('/api/preview-colors', methods=['POST'])
def preview_colors():
    """Preview colors from uploaded logo"""
    if 'logo' not in request.files:
        return jsonify(default_colors())
    logo_file = request.files['logo']
    if not logo_file.filename:
        return jsonify(default_colors())
    
    temp_path = UPLOAD_DIR / f"temp_{uuid.uuid4().hex[:8]}{Path(logo_file.filename).suffix}"
    logo_file.save(str(temp_path))
    colors = extract_colors_from_logo(str(temp_path))
    try: os.unlink(str(temp_path))
    except: pass
    return jsonify(colors)

# ── Main Page ─────────────────────────────────────────────────
LANDING_HTML = """<!DOCTYPE html>
<html lang="en"><head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>ProposalSnap — AI Presentation Maker</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&family=JetBrains+Mono:wght@400;500&display=swap" rel="stylesheet">
<style>
:root{--bg:#0B0F1A;--surface:#131829;--border:rgba(255,255,255,0.08);--text:#F0F0F5;
--text2:#8B8FA3;--accent:#6C5CE7;--accent2:#A78BFA;--green:#00D2A0;--red:#FF6B6B;
--radius:14px;--font:'Inter',sans-serif}
*{margin:0;padding:0;box-sizing:border-box}
body{font-family:var(--font);background:var(--bg);color:var(--text);min-height:100vh}
a{text-decoration:none;color:inherit}
.btn-hero{display:inline-block;padding:14px 32px;border-radius:10px;font-size:15px;font-weight:700;margin:0 6px 8px;transition:all .2s;cursor:pointer}
.btn-hero:hover{transform:translateY(-2px)}
.btn-fill{background:linear-gradient(135deg,#6C5CE7,#5A4BD1);color:#fff !important;box-shadow:0 4px 20px rgba(108,92,231,.2)}
.btn-fill:hover{box-shadow:0 8px 30px rgba(108,92,231,.4)}
.btn-outline{background:transparent;color:var(--text) !important;border:1.5px solid rgba(255,255,255,.15)}
.btn-outline:hover{border-color:var(--accent);color:var(--accent) !important}
</style>
</head>
<body>
<section style="min-height:80vh;display:flex;align-items:center;justify-content:center;text-align:center;padding:60px 24px 40px;position:relative;overflow:hidden">
<div style="position:absolute;width:600px;height:600px;background:radial-gradient(circle,rgba(108,92,231,.12),transparent 70%);top:-100px;right:-100px;pointer-events:none;z-index:0"></div>
<div style="max-width:700px;position:relative;z-index:1">
<div style="display:inline-flex;align-items:center;gap:8px;background:rgba(108,92,231,.1);border:1px solid rgba(108,92,231,.2);border-radius:20px;padding:6px 16px;font-size:12px;font-weight:700;color:var(--accent);margin-bottom:24px">✦ AI-Powered Presentations</div>
<h1 style="font-size:clamp(32px,5vw,48px);font-weight:800;line-height:1.15;margin-bottom:16px;color:#fff">Describe your idea.<br><span style="background:linear-gradient(135deg,#6C5CE7,#00D2A0);-webkit-background-clip:text;-webkit-text-fill-color:transparent">Get a pitch deck.</span></h1>
<p style="font-size:17px;color:var(--text2);line-height:1.7;margin-bottom:28px;max-width:560px;margin:0 auto 28px">Create presentations from scratch or polish existing decks with AI. Smart layouts, branded colors, structured slides. Ready to present in 30 seconds.</p>
<div>
<a href="/create" class="btn-hero btn-fill">Create Presentation →</a>
<a href="/login" class="btn-hero btn-outline">Sign In</a>
<a href="/register" class="btn-hero btn-outline">Create Account</a>
<a href="#features" class="btn-hero btn-outline">See Features</a>
<a href="/demo" class="btn-hero btn-outline">View Demo Gallery</a>
</div>
</div>
</section>

<section id="features" style="padding:60px 24px;max-width:1000px;margin:0 auto">
<div style="font-size:12px;font-weight:800;text-transform:uppercase;letter-spacing:2px;color:var(--accent);margin-bottom:12px;text-align:center">Features</div>
<div style="font-size:28px;font-weight:800;color:#fff;text-align:center;margin-bottom:36px">Professional decks without the design work</div>
<div style="display:grid;grid-template-columns:repeat(auto-fit,minmax(280px,1fr));gap:16px">
<div style="background:var(--surface);border:1px solid var(--border);border-radius:14px;padding:22px"><div style="font-size:26px;margin-bottom:10px">🤖</div><div style="font-size:15px;font-weight:700;color:#fff;margin-bottom:6px">AI Slide Generation</div><div style="font-size:13px;color:var(--text2);line-height:1.6">Describe your project — AI creates 15+ layout types: timelines, stats, comparisons, icon grids, pricing tables, and more.</div></div>
<div style="background:var(--surface);border:1px solid var(--border);border-radius:14px;padding:22px"><div style="font-size:26px;margin-bottom:10px">✨</div><div style="font-size:15px;font-weight:700;color:#fff;margin-bottom:6px">Polish Existing Decks</div><div style="font-size:13px;color:var(--text2);line-height:1.6">Upload any .pptx — AI improves tone, restructures slides, adds visuals. Quick presets: Make Concise, Investor-Ready, Executive Summary.</div></div>
<div style="background:var(--surface);border:1px solid var(--border);border-radius:14px;padding:22px"><div style="font-size:26px;margin-bottom:10px">👥</div><div style="font-size:15px;font-weight:700;color:#fff;margin-bottom:6px">Audience Versioning</div><div style="font-size:13px;color:var(--text2);line-height:1.6">One deck, three versions: 6-slide exec summary, 14-slide team detail, and 10-slide investor pitch — generated simultaneously.</div></div>
<div style="background:var(--surface);border:1px solid var(--border);border-radius:14px;padding:22px"><div style="font-size:26px;margin-bottom:10px">🎭</div><div style="font-size:15px;font-weight:700;color:#fff;margin-bottom:6px">Style Transfer</div><div style="font-size:13px;color:var(--text2);line-height:1.6">Upload a reference deck you love + your content deck. AI copies the style, structure, and rhythm while keeping your data.</div></div>
<div style="background:var(--surface);border:1px solid var(--border);border-radius:14px;padding:22px"><div style="font-size:26px;margin-bottom:10px">🔗</div><div style="font-size:15px;font-weight:700;color:#fff;margin-bottom:6px">Merge & Split</div><div style="font-size:13px;color:var(--text2);line-height:1.6">Combine slides from multiple team members into one cohesive deck. Or split a monster deck into an exec summary + appendix.</div></div>
<div style="background:var(--surface);border:1px solid var(--border);border-radius:14px;padding:22px"><div style="font-size:26px;margin-bottom:10px">🎨</div><div style="font-size:15px;font-weight:700;color:#fff;margin-bottom:6px">Brand-Locked Output</div><div style="font-size:13px;color:var(--text2);line-height:1.6">Save your logo once — every deck comes out in your exact brand colors and fonts. Consistent branding across all presentations.</div></div>
</div>
</section>

<section style="padding:40px 24px 60px;max-width:800px;margin:0 auto">
<div style="font-size:12px;font-weight:800;text-transform:uppercase;letter-spacing:2px;color:#00D2A0;margin-bottom:12px;text-align:center">How it works</div>
<div style="font-size:28px;font-weight:800;color:#fff;text-align:center;margin-bottom:32px">Three inputs. One deck.</div>
<div style="display:flex;flex-direction:column;gap:14px">
<div style="display:flex;gap:14px;align-items:flex-start;background:var(--surface);border:1px solid var(--border);border-radius:14px;padding:18px"><div style="min-width:34px;height:34px;border-radius:10px;background:linear-gradient(135deg,#6C5CE7,#5A4BD1);color:#fff;display:flex;align-items:center;justify-content:center;font-weight:800;font-size:14px;flex-shrink:0">1</div><div><div style="font-size:14px;font-weight:700;color:#fff;margin-bottom:3px">Describe your project</div><div style="font-size:13px;color:var(--text2);line-height:1.5">Write a brief — be as detailed or vague as you want. AI fills in the rest.</div></div></div>
<div style="display:flex;gap:14px;align-items:flex-start;background:var(--surface);border:1px solid var(--border);border-radius:14px;padding:18px"><div style="min-width:34px;height:34px;border-radius:10px;background:linear-gradient(135deg,#6C5CE7,#5A4BD1);color:#fff;display:flex;align-items:center;justify-content:center;font-weight:800;font-size:14px;flex-shrink:0">2</div><div><div style="font-size:14px;font-weight:700;color:#fff;margin-bottom:3px">Pick slides & colors</div><div style="font-size:13px;color:var(--text2);line-height:1.5">Choose slide count (5–20), select a color theme or upload your logo for auto-extraction.</div></div></div>
<div style="display:flex;gap:14px;align-items:flex-start;background:var(--surface);border:1px solid var(--border);border-radius:14px;padding:18px"><div style="min-width:34px;height:34px;border-radius:10px;background:linear-gradient(135deg,#6C5CE7,#5A4BD1);color:#fff;display:flex;align-items:center;justify-content:center;font-weight:800;font-size:14px;flex-shrink:0">3</div><div><div style="font-size:14px;font-weight:700;color:#fff;margin-bottom:3px">Generate & download</div><div style="font-size:13px;color:var(--text2);line-height:1.5">AI creates the full deck. Download the .pptx and present or edit freely.</div></div></div>
</div>
</section>

<section style="padding:40px 24px 80px;text-align:center">
<a href="/create" style="display:inline-block;padding:18px 48px;background:linear-gradient(135deg,#6C5CE7,#5A4BD1);color:#fff !important;border-radius:12px;font-size:17px;font-weight:700;transition:.2s;box-shadow:0 4px 20px rgba(108,92,231,.3);margin:0 8px 10px">Create Presentation →</a>
<a href="/register" style="display:inline-block;padding:18px 48px;background:transparent;color:var(--text) !important;border:1.5px solid rgba(255,255,255,.15);border-radius:12px;font-size:17px;font-weight:600;margin:0 8px 10px">Create Free Account</a>

<div style="margin-top:20px;font-size:12px;color:var(--text2)">Part of <a href="https://snapsuite.up.railway.app" style="color:#A78BFA !important">SnapSuite</a> — 6 apps for your entire business</div>
</section>
</body></html>"""






MAIN_HTML = """<!DOCTYPE html>
<html lang="en"><head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>ProposalSnap — AI Presentation Maker</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&family=JetBrains+Mono:wght@400;500&display=swap" rel="stylesheet">
<style>
:root{--bg:#0B0F1A;--surface:#131829;--border:rgba(255,255,255,0.08);--text:#F0F0F5;
--text2:#8B8FA3;--accent:#6C5CE7;--accent2:#A78BFA;--green:#00D2A0;--red:#FF6B6B;
--radius:14px;--font:'Inter',sans-serif}
*{margin:0;padding:0;box-sizing:border-box}
body{font-family:var(--font);background:var(--bg);color:var(--text);min-height:100vh}
.container{max-width:900px;margin:0 auto;padding:24px 16px}
h1{font-size:32px;font-weight:700;background:linear-gradient(135deg,#6C5CE7,#00D2A0);
-webkit-background-clip:text;-webkit-text-fill-color:transparent;margin-bottom:4px}
.subtitle{color:var(--text2);font-size:15px;margin-bottom:32px}
.card{background:var(--surface);border:1px solid var(--border);border-radius:var(--radius);padding:24px;margin-bottom:16px}
.card h3{font-size:16px;font-weight:600;margin-bottom:16px;color:var(--accent2)}
label{display:block;font-size:13px;color:var(--text2);margin-bottom:6px;font-weight:500}
input[type="text"],textarea,select{width:100%;padding:12px 16px;background:var(--bg);
border:1px solid var(--border);border-radius:10px;color:var(--text);font-family:var(--font);
font-size:14px;outline:none;transition:border-color 0.2s}
input:focus,textarea:focus,select:focus{border-color:var(--accent)}
textarea{resize:vertical;min-height:120px;line-height:1.6}
select{cursor:pointer;appearance:none;background-image:url("data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' width='12' height='12' viewBox='0 0 12 12'%3E%3Cpath d='M6 8L1 3h10z' fill='%238B8FA3'/%3E%3C/svg%3E");
background-repeat:no-repeat;background-position:right 12px center}
.row{display:grid;grid-template-columns:1fr 1fr;gap:16px}
@media(max-width:600px){.row{grid-template-columns:1fr}}
.row3{display:grid;grid-template-columns:1fr 1fr 1fr 1fr;gap:16px}
@media(max-width:600px){.row3{grid-template-columns:1fr 1fr}}
.btn{padding:14px 28px;border-radius:10px;font-family:var(--font);font-size:15px;font-weight:600;
cursor:pointer;border:none;transition:all 0.2s}
.btn-primary{background:linear-gradient(135deg,#6C5CE7,#5A4BD1);color:white;width:100%}
.btn-primary:hover{transform:translateY(-1px);box-shadow:0 4px 15px rgba(108,92,231,0.4)}
.btn-primary:disabled{opacity:0.5;transform:none;cursor:not-allowed}
.file-upload{border:2px dashed var(--border);border-radius:10px;padding:24px;text-align:center;
cursor:pointer;transition:all 0.2s}
.file-upload:hover{border-color:var(--accent);background:rgba(108,92,231,0.05)}
.file-upload.has-file{border-color:var(--green);background:rgba(0,210,160,0.05)}
.file-upload input{display:none}
.color-preview{display:flex;gap:8px;margin-top:12px;align-items:center}
.color-dot{width:32px;height:32px;border-radius:8px;border:2px solid var(--border)}
.color-label{font-size:11px;color:var(--text2);text-align:center;margin-top:2px}
.status{padding:16px;border-radius:10px;margin-top:16px;display:none}
.status.loading{display:block;background:rgba(108,92,231,0.1);color:var(--accent2)}
.status.success{display:block;background:rgba(0,210,160,0.1);color:var(--green)}
.status.error{display:block;background:rgba(255,107,107,0.1);color:var(--red)}
.download-btn{display:inline-flex;align-items:center;gap:8px;padding:12px 24px;
background:var(--green);color:var(--bg);border-radius:10px;text-decoration:none;
font-weight:600;font-size:14px;margin-top:12px;transition:all 0.2s}
.download-btn:hover{transform:translateY(-1px);box-shadow:0 4px 15px rgba(0,210,160,0.3)}
.spinner{display:inline-block;width:16px;height:16px;border:2px solid rgba(255,255,255,0.3);
border-top-color:var(--accent2);border-radius:50%;animation:spin 0.8s linear infinite}
@keyframes spin{to{transform:rotate(360deg)}}
.examples{display:flex;gap:8px;margin-top:8px;flex-wrap:wrap}
.example-tag{font-size:11px;padding:4px 10px;background:rgba(108,92,231,0.1);color:var(--accent2);
border-radius:20px;cursor:pointer;border:1px solid transparent;transition:all 0.2s}
.example-tag:hover{border-color:var(--accent);background:rgba(108,92,231,0.2)}
.footer{text-align:center;padding:32px 0;color:var(--text2);font-size:13px}
.counter{text-align:right;font-size:12px;color:var(--text2);margin-top:4px}
</style>
</head>
<body>
<div class="container">
<div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:8px">
<a href="/" style="text-decoration:none;color:inherit"><h1>ProposalSnap</h1></a>
<a href="https://snapsuite.up.railway.app" target="_blank" style="font-size:12px;color:#8B95B0;text-decoration:none;padding:6px 12px;border:1px solid #2A3148;border-radius:6px;font-weight:600;font-family:'DM Sans',sans-serif">← SnapSuite</a>
<a href="/logout" style="font-size:12px;color:#FF6B6B;text-decoration:none;padding:6px 12px;border:1px solid rgba(255,107,107,.3);border-radius:6px;font-weight:600;font-family:'DM Sans',sans-serif">Sign Out</a>
{% if is_admin %}<a href="/admin" style="font-size:12px;color:#A78BFA;text-decoration:none;padding:6px 12px;border:1px solid rgba(167,139,250,.3);border-radius:6px;font-weight:600;font-family:'DM Sans',sans-serif">Admin</a>{% endif %}

<div style="position:relative;display:inline-block"><button onclick="this.nextElementSibling.style.display=this.nextElementSibling.style.display==='block'?'none':'block'" style="font-size:14px;background:none;border:1px solid #2A3148;border-radius:6px;padding:5px 10px;color:#8B95B0;cursor:pointer;font-family:'DM Sans',sans-serif" title="Switch App">⊞</button><div style="display:none;position:absolute;right:0;top:32px;background:#141926;border:1px solid #2A3148;border-radius:10px;padding:8px;min-width:180px;z-index:200;box-shadow:0 8px 30px rgba(0,0,0,.5)"><a href="https://invoicesnap.up.railway.app" style="display:block;padding:8px 12px;color:#E8ECF4;text-decoration:none;border-radius:6px;font-size:13px;font-weight:500;font-family:'DM Sans',sans-serif" onmouseover="this.style.background='#2A3148'" onmouseout="this.style.background='none'">📄 InvoiceSnap</a><a href="https://contractsnap-app.up.railway.app" style="display:block;padding:8px 12px;color:#E8ECF4;text-decoration:none;border-radius:6px;font-size:13px;font-weight:500;font-family:'DM Sans',sans-serif" onmouseover="this.style.background='#2A3148'" onmouseout="this.style.background='none'">📋 ContractSnap</a><a href="https://expensesnap.up.railway.app" style="display:block;padding:8px 12px;color:#E8ECF4;text-decoration:none;border-radius:6px;font-size:13px;font-weight:500;font-family:'DM Sans',sans-serif" onmouseover="this.style.background='#2A3148'" onmouseout="this.style.background='none'">📸 ExpenseSnap</a><a href="https://payslipsnap.up.railway.app" style="display:block;padding:8px 12px;color:#E8ECF4;text-decoration:none;border-radius:6px;font-size:13px;font-weight:500;font-family:'DM Sans',sans-serif" onmouseover="this.style.background='#2A3148'" onmouseout="this.style.background='none'">💰 PayslipSnap</a></div></div>
</div>
<p class="subtitle">Create new presentations or polish existing ones with AI</p>

<div style="display:flex;gap:0;margin-bottom:24px;background:var(--surface);border:1px solid var(--border);border-radius:12px;padding:4px;flex-wrap:wrap">
<button onclick="switchMode('create')" id="tabCreate" class="modetab active" style="flex:1;min-width:100px;padding:10px 14px;border:none;border-radius:10px;font-family:var(--font);font-size:13px;font-weight:600;cursor:pointer;transition:all .2s;background:linear-gradient(135deg,#6C5CE7,#5A4BD1);color:#fff">⚡ Create</button>
<button onclick="switchMode('polish')" id="tabPolish" class="modetab" style="flex:1;min-width:100px;padding:10px 14px;border:none;border-radius:10px;font-family:var(--font);font-size:13px;font-weight:600;cursor:pointer;transition:all .2s;background:transparent;color:var(--text2)">✨ Polish</button>
<button onclick="switchMode('version')" id="tabVersion" class="modetab" style="flex:1;min-width:100px;padding:10px 14px;border:none;border-radius:10px;font-family:var(--font);font-size:13px;font-weight:600;cursor:pointer;transition:all .2s;background:transparent;color:var(--text2)">👥 Versions</button>
<button onclick="switchMode('style')" id="tabStyle" class="modetab" style="flex:1;min-width:100px;padding:10px 14px;border:none;border-radius:10px;font-family:var(--font);font-size:13px;font-weight:600;cursor:pointer;transition:all .2s;background:transparent;color:var(--text2)">🎭 Style Transfer</button>
<button onclick="switchMode('merge')" id="tabMerge" class="modetab" style="flex:1;min-width:100px;padding:10px 14px;border:none;border-radius:10px;font-family:var(--font);font-size:13px;font-weight:600;cursor:pointer;transition:all .2s;background:transparent;color:var(--text2)">🔗 Merge & Split</button>
</div>

<!-- ═══ CREATE MODE ═══ -->
<div id="modeCreate">
<div class="card">
<h3>📋 Presentation Details</h3>
<div class="row" style="margin-bottom:16px">
<div><label>Client Name *</label>
<input type="text" id="clientName" placeholder="e.g. Zurich Insurance"></div>
<div><label>Your Company Name</label>
<input type="text" id="companyName" placeholder="e.g. Shakty.AI"></div>
</div>
<div class="row3" style="margin-bottom:16px">
<div><label>Presentation Type</label>
<select id="presType">
<option>Corporate Proposal</option>
<option>Sales Pitch Deck</option>
<option>Training Deck</option>
<option>Project Report</option>
<option>Business Plan</option>
<option>Consulting Engagement</option>
</select></div>
<div><label>Tone</label>
<select id="tone">
<option>Corporate</option>
<option>Creative</option>
<option>Minimal</option>
<option>Bold</option>
<option>Friendly</option>
</select></div>
<div><label>Font Style</label>
<select id="fontStyle">
<option value="aptos">Aptos · Clean Modern</option>
<option value="georgia">Georgia + Calibri · Classic</option>
<option value="arial">Arial Black + Arial · Bold</option>
<option value="trebuchet">Trebuchet + Calibri · Creative</option>
<option value="palatino">Palatino + Garamond · Elegant</option>
<option value="cambria">Cambria + Calibri · Traditional</option>
</select></div>
<div><label>Number of Slides</label>
<select id="numSlides">
<option>1</option>
<option>3</option>
<option>6</option>
<option>8</option>
<option>10</option>
<option selected>12</option>
<option>14</option>
<option>16</option>
</select></div>
</div>
</div>

<div class="card">
<h3>🎨 Logo & Branding</h3>
<p style="color:var(--text2);font-size:13px;margin-bottom:12px">Upload your logo — we'll extract colors to style the entire presentation.</p>
<div class="file-upload" id="logoDropzone" onclick="document.getElementById('logoInput').click()">
<input type="file" id="logoInput" accept=".png,.jpg,.jpeg,.webp,.svg" onchange="handleLogo(this)">
<div id="logoText">📎 Click to upload logo (PNG, JPG, SVG)</div>
</div>
<div class="color-preview" id="colorPreview" style="display:none;">
<div><div class="color-dot" id="cprimary"></div><div class="color-label">Primary</div></div>
<div><div class="color-dot" id="csecondary"></div><div class="color-label">Secondary</div></div>
<div><div class="color-dot" id="caccent"></div><div class="color-label">Accent</div></div>
<div><div class="color-dot" id="cdark"></div><div class="color-label">Dark</div></div>
<div style="flex:1;font-size:12px;color:var(--text2);padding-left:8px">Colors extracted from your logo</div>
</div>
</div>

<div class="card">
<h3>💡 Key Points & Content *</h3>
<p style="color:var(--text2);font-size:13px;margin-bottom:12px">Describe what the presentation should cover. The more detail, the better the output.</p>
<textarea id="keyPoints" placeholder="Example:
- We offer AI-powered expense management for small businesses
- Key differentiator: Receipt scanning with high accuracy using AI  
- Target market: Consulting firms with 5-50 employees
- Implementation: 2-week setup, full team training included
- Benefits: Save significant time on expense reporting per employee
- Security: HTTPS, data isolation, daily backups"></textarea>
<div class="counter"><span id="charCount">0</span> characters</div>
<div class="examples">
<span class="example-tag" onclick="fillExample('proposal')">💼 Proposal</span>
<span class="example-tag" onclick="fillExample('pitch')">🚀 Pitch Deck</span>
<span class="example-tag" onclick="fillExample('training')">📚 Training</span>
<span class="example-tag" onclick="fillExample('report')">📊 Report</span>
</div>
</div>

<button class="btn btn-primary" id="generateBtn" onclick="generate()">
⚡ Generate Presentation
</button>

<div class="status" id="status"></div>
</div>

<!-- ═══ POLISH MODE ═══ -->
<div id="modePolish" style="display:none">
<div class="card">
<h3>📤 Upload Your Deck</h3>
<p style="color:var(--text2);font-size:13px;margin-bottom:12px">Upload an existing .pptx file — AI will extract the content and polish it based on your instructions.</p>
<div class="file-upload" id="pptxDropzone" onclick="document.getElementById('pptxInput').click()">
<input type="file" id="pptxInput" accept=".pptx" onchange="handlePptx(this)">
<div id="pptxText">📎 Click to upload your PPTX file</div>
</div>
</div>

<div class="card">
<h3>✏️ Polishing Instructions *</h3>
<p style="color:var(--text2);font-size:13px;margin-bottom:12px">Tell the AI what to improve, restructure, or change. Be specific for best results.</p>
<textarea id="polishInstructions" placeholder="Examples:
- Make it more concise and executive-friendly
- Add a timeline slide and a comparison table
- Change the tone to be more persuasive
- Improve the visuals — use more charts and stats
- Restructure: combine slides 3 & 4, expand the key benefits section
- Add a team slide and a competitive advantage comparison
- Make all bullets more action-oriented"></textarea>
<div class="counter"><span id="polishCharCount">0</span> characters</div>
<div class="examples">
<span class="example-tag" onclick="fillPolishExample('concise')">✂️ Make Concise</span>
<span class="example-tag" onclick="fillPolishExample('visual')">📊 More Visual</span>
<span class="example-tag" onclick="fillPolishExample('investor')">💰 Investor-Ready</span>
<span class="example-tag" onclick="fillPolishExample('executive')">👔 Executive Summary</span>
</div>
</div>

<div class="card">
<h3>🎨 Tone & Style</h3>
<div class="row" style="margin-bottom:16px">
<div><label>Client / Deck Name</label>
<input type="text" id="polishClientName" placeholder="e.g. Q4 Strategy Review"></div>
<div><label>Company Name</label>
<input type="text" id="polishCompanyName" placeholder="e.g. Shakty.AI"></div>
</div>
<div class="row3">
<div><label>Tone</label>
<select id="polishTone">
<option>Corporate</option>
<option>Creative</option>
<option>Minimal</option>
<option>Bold</option>
<option>Friendly</option>
</select></div>
<div><label>Font Style</label>
<select id="polishFontStyle">
<option value="aptos">Aptos · Clean Modern</option>
<option value="georgia">Georgia + Calibri · Classic</option>
<option value="arial">Arial Black + Arial · Bold</option>
<option value="trebuchet">Trebuchet + Calibri · Creative</option>
<option value="palatino">Palatino + Garamond · Elegant</option>
<option value="cambria">Cambria + Calibri · Traditional</option>
</select></div>
<div><label>Output Slides</label>
<select id="polishNumSlides">
<option value="">Same as original</option>
<option>1</option>
<option>3</option>
<option>6</option>
<option>8</option>
<option>10</option>
<option>12</option>
<option>14</option>
<option>16</option>
</select></div>
<div><label>Logo (optional)</label>
<div class="file-upload" style="padding:10px" onclick="document.getElementById('polishLogoInput').click()">
<input type="file" id="polishLogoInput" accept=".png,.jpg,.jpeg,.webp,.svg" onchange="handlePolishLogo(this)">
<div id="polishLogoText" style="font-size:12px">📎 Upload logo</div>
</div></div>
</div>
</div>

<button class="btn btn-primary" id="polishBtn" onclick="polishDeck()">
✨ Polish Presentation
</button>

<div class="status" id="polishStatus"></div>
</div>

<!-- ═══ AUDIENCE VERSIONS MODE ═══ -->
<div id="modeVersion" style="display:none">
<div class="card">
<h3>📤 Upload Your Deck</h3>
<p style="color:var(--text2);font-size:13px;margin-bottom:12px">Upload one deck and get 3 tailored versions — for executives, team, and investors.</p>
<div class="file-upload" id="versionDropzone" onclick="document.getElementById('versionInput').click()">
<input type="file" id="versionInput" accept=".pptx" onchange="handleFileUpload(this,'versionDropzone','versionText')">
<div id="versionText">📎 Click to upload your PPTX file</div>
</div>
</div>

<div class="card">
<h3>👥 Select Audience Versions</h3>
<p style="color:var(--text2);font-size:13px;margin-bottom:14px">Each version restructures your content for a different audience.</p>
<div style="display:flex;flex-direction:column;gap:10px">
<label style="display:flex;align-items:flex-start;gap:10px;padding:14px;background:var(--bg);border:1px solid var(--border);border-radius:10px;cursor:pointer;font-size:13px">
<input type="checkbox" id="vExec" checked style="margin-top:2px"> <div><strong style="color:#fff">👔 Executive Summary</strong><br><span style="color:var(--text2)">6 slides max. Lead with conclusion. Stats over bullets. Answerable in 3 minutes.</span></div></label>
<label style="display:flex;align-items:flex-start;gap:10px;padding:14px;background:var(--bg);border:1px solid var(--border);border-radius:10px;cursor:pointer;font-size:13px">
<input type="checkbox" id="vDetailed" checked style="margin-top:2px"> <div><strong style="color:#fff">📋 Detailed (Team)</strong><br><span style="color:var(--text2)">12-16 slides. Full context, process flows, timelines, and supporting data.</span></div></label>
<label style="display:flex;align-items:flex-start;gap:10px;padding:14px;background:var(--bg);border:1px solid var(--border);border-radius:10px;cursor:pointer;font-size:13px">
<input type="checkbox" id="vInvestor" checked style="margin-top:2px"> <div><strong style="color:#fff">💰 Investor Pitch</strong><br><span style="color:var(--text2)">8-10 slides. Problem → Solution → Market → Traction → Ask. Numbers prominent.</span></div></label>
</div>
</div>

<div class="card">
<h3>🎨 Settings</h3>
<div class="row">
<div><label>Deck Name</label><input type="text" id="vClientName" placeholder="e.g. Q4 Strategy Review"></div>
<div><label>Company Name</label><input type="text" id="vCompanyName" placeholder="e.g. Shakty.AI"></div>
</div>
<div class="row" style="margin-top:14px">
<div><label>Font Style</label><select id="vFontStyle"><option value="aptos">Aptos · Clean Modern</option><option value="georgia">Georgia + Calibri · Classic</option><option value="arial">Arial Black · Bold</option><option value="trebuchet">Trebuchet · Creative</option><option value="palatino">Palatino · Elegant</option><option value="cambria">Cambria · Traditional</option></select></div>
<div><label>Logo (optional)</label><div class="file-upload" style="padding:10px" onclick="document.getElementById('vLogoInput').click()"><input type="file" id="vLogoInput" accept=".png,.jpg,.jpeg,.webp,.svg" onchange="this.parentElement.querySelector('div').textContent='✓ '+this.files[0].name"><div style="font-size:12px">📎 Upload logo</div></div></div>
</div>
</div>

<button class="btn btn-primary" id="versionBtn" onclick="generateVersions()">
👥 Generate All Versions
</button>
<div class="status" id="versionStatus"></div>
</div>

<!-- ═══ STYLE TRANSFER MODE ═══ -->
<div id="modeStyle" style="display:none">
<div class="card">
<h3>📄 Your Content Deck</h3>
<p style="color:var(--text2);font-size:13px;margin-bottom:12px">Upload the deck with the content you want to keep.</p>
<div class="file-upload" id="contentDropzone" onclick="document.getElementById('contentInput').click()">
<input type="file" id="contentInput" accept=".pptx" onchange="handleFileUpload(this,'contentDropzone','contentText')">
<div id="contentText">📎 Upload your content PPTX</div>
</div>
</div>

<div class="card">
<h3>🎭 Reference / Style Deck</h3>
<p style="color:var(--text2);font-size:13px;margin-bottom:12px">Upload a deck whose style you want to copy — structure, tone, rhythm, level of detail.</p>
<div class="file-upload" id="refDropzone" onclick="document.getElementById('refInput').click()">
<input type="file" id="refInput" accept=".pptx" onchange="handleFileUpload(this,'refDropzone','refText')">
<div id="refText">📎 Upload reference/style PPTX</div>
</div>
</div>

<div class="card">
<h3>✏️ Additional Instructions (optional)</h3>
<textarea id="styleInstructions" rows="3" placeholder="e.g. Match the McKinsey style but keep our data. Make it more visual."></textarea>
</div>

<div class="card">
<h3>🎨 Settings</h3>
<div class="row">
<div><label>Deck Name</label><input type="text" id="sClientName" placeholder="e.g. Strategy Proposal"></div>
<div><label>Company Name</label><input type="text" id="sCompanyName" placeholder="e.g. Shakty.AI"></div>
</div>
<div class="row" style="margin-top:14px">
<div><label>Font Style</label><select id="sFontStyle"><option value="aptos">Aptos · Clean Modern</option><option value="georgia">Georgia + Calibri · Classic</option><option value="arial">Arial Black · Bold</option><option value="trebuchet">Trebuchet · Creative</option><option value="palatino">Palatino · Elegant</option><option value="cambria">Cambria · Traditional</option></select></div>
<div><label>Logo (optional)</label><div class="file-upload" style="padding:10px" onclick="document.getElementById('sLogoInput').click()"><input type="file" id="sLogoInput" accept=".png,.jpg,.jpeg,.webp,.svg" onchange="this.parentElement.querySelector('div').textContent='✓ '+this.files[0].name"><div style="font-size:12px">📎 Upload logo</div></div></div>
</div>
</div>

<button class="btn btn-primary" id="styleBtn" onclick="doStyleTransfer()">
🎭 Apply Style Transfer
</button>
<div class="status" id="styleStatus"></div>
</div>

<!-- ═══ MERGE & SPLIT MODE ═══ -->
<div id="modeMerge" style="display:none">
<div class="card">
<h3>🔗 Merge Multiple Decks</h3>
<p style="color:var(--text2);font-size:13px;margin-bottom:12px">Upload 2+ PPTX files — AI will combine them into one cohesive presentation.</p>
<div class="file-upload" id="mergeDropzone" onclick="document.getElementById('mergeInput').click()">
<input type="file" id="mergeInput" accept=".pptx" multiple onchange="handleMultiFile(this)">
<div id="mergeText">📎 Click to upload multiple PPTX files</div>
</div>
<div style="margin-top:12px"><label>Merge Instructions (optional)</label>
<textarea id="mergeInstructions" rows="2" placeholder="e.g. Combine into a single narrative. Remove duplicates. Prioritize team A's data slides."></textarea></div>
</div>

<div class="card">
<h3>✂️ Split a Deck</h3>
<p style="color:var(--text2);font-size:13px;margin-bottom:12px">Upload one large deck — get an Executive Summary + Full Detail version.</p>
<div class="file-upload" id="splitDropzone" onclick="document.getElementById('splitInput').click()">
<input type="file" id="splitInput" accept=".pptx" onchange="handleFileUpload(this,'splitDropzone','splitText')">
<div id="splitText">📎 Upload PPTX to split</div>
</div>
</div>

<div class="card">
<h3>🎨 Settings</h3>
<div class="row">
<div><label>Deck Name</label><input type="text" id="mClientName" placeholder="e.g. Q4 Report"></div>
<div><label>Company Name</label><input type="text" id="mCompanyName" placeholder="e.g. Shakty.AI"></div>
</div>
<div class="row" style="margin-top:14px">
<div><label>Font Style</label><select id="mFontStyle"><option value="aptos">Aptos · Clean Modern</option><option value="georgia">Georgia + Calibri · Classic</option><option value="arial">Arial Black · Bold</option><option value="trebuchet">Trebuchet · Creative</option><option value="palatino">Palatino · Elegant</option><option value="cambria">Cambria · Traditional</option></select></div>
<div><label>Tone</label><select id="mTone"><option>Corporate</option><option>Creative</option><option>Minimal</option><option>Bold</option><option>Friendly</option></select></div>
</div>
</div>

<div style="display:flex;gap:12px;flex-wrap:wrap">
<button class="btn btn-primary" id="mergeBtn" onclick="doMerge()" style="flex:1;min-width:200px">
🔗 Merge Decks
</button>
<button class="btn btn-primary" id="splitBtn" onclick="doSplit()" style="flex:1;min-width:200px">
✂️ Split Deck
</button>
</div>
<div class="status" id="mergeStatus"></div>
</div>

<div class="footer">
ProposalSnap · Built with Claude AI · Powered by Shakty.AI
</div>
</div>

<script>
const keyPointsEl = document.getElementById('keyPoints');
keyPointsEl.addEventListener('input', () => {
  document.getElementById('charCount').textContent = keyPointsEl.value.length;
});

async function handleLogo(input) {
  if (!input.files.length) return;
  const file = input.files[0];
  document.getElementById('logoText').textContent = `✓ ${file.name}`;
  document.getElementById('logoDropzone').classList.add('has-file');
  
  // Preview colors
  const formData = new FormData();
  formData.append('logo', file);
  try {
    const res = await fetch('/api/preview-colors', { method: 'POST', body: formData });
    const colors = await res.json();
    document.getElementById('cprimary').style.background = '#' + colors.primary;
    document.getElementById('csecondary').style.background = '#' + colors.secondary;
    document.getElementById('caccent').style.background = '#' + colors.accent;
    document.getElementById('cdark').style.background = '#' + colors.dark;
    document.getElementById('colorPreview').style.display = 'flex';
  } catch(e) { console.error(e); }
}

function fillExample(type) {
  const examples = {
    proposal: `Client wants an AI-powered expense management solution
- Receipt scanning with high accuracy using AI
- Multi-company support with role-based access
- Real-time currency conversion across multiple currencies
- Key differentiator: works on phone camera, no app install needed
- Target: consulting firms with small to mid-size teams
- Flexible pricing tiers: Starter, Business, and Pro
- Implementation: 2-week setup, full team training included
- Benefits: Significant time savings on expense reporting per employee
- Security: HTTPS, data isolation, PostgreSQL with daily backups`,
    pitch: `We are building the next generation of AI productivity tools
- Problem: Small businesses waste significant time on admin tasks
- Solution: AI agents that automate receipts, invoices, and reporting
- Large and growing global expense management market
- Early traction with paying clients and growing revenue
- Technology: AI-powered extraction, cloud-hosted infrastructure
- Team: Experienced founders with decades of combined experience
- Seeking seed funding for product development and growth
- Use of funds: Product development, Sales, Infrastructure`,
    training: `AI Training Program for Finance Team
- Module 1: Introduction to AI in Finance (Copilot, Claude, ChatGPT)
- Module 2: Prompt Engineering for Financial Analysis
- Module 3: Building AI Agents with Copilot Studio
- Module 4: Automating Expense Reporting and Invoice Processing
- Module 5: AI-Powered Financial Dashboards
- Duration: 6 weeks, 2 sessions per week
- Target audience: Finance managers and analysts
- Expected outcomes: Major reduction in manual tasks
- Certification provided upon completion`,
    report: `Q4 Financial Performance Summary
- Revenue growth year-over-year with strong momentum
- New enterprise accounts acquired during the quarter
- Customer retention above industry average
- Key wins: Major contracts signed with leading companies
- Challenges: Exchange rate fluctuations, supply chain delays
- Cost optimization: Significant reduction in operational costs
- Team growth during the quarter
- Next quarter outlook: Strong pipeline, targeting continued growth
- Strategic priorities: AI implementation, compliance, market expansion`
  };
  keyPointsEl.value = examples[type] || '';
  document.getElementById('charCount').textContent = keyPointsEl.value.length;
}

async function generate() {
  const clientName = document.getElementById('clientName').value.trim();
  const companyName = document.getElementById('companyName').value.trim();
  const presType = document.getElementById('presType').value;
  const tone = document.getElementById('tone').value;
  const fontStyle = document.getElementById('fontStyle').value;
  const keyPoints = document.getElementById('keyPoints').value.trim();
  const numSlides = document.getElementById('numSlides').value;
  
  if (!clientName) { showStatus('Please enter a client name', 'error'); return; }
  if (!keyPoints) { showStatus('Please enter key points for the presentation', 'error'); return; }
  
  const btn = document.getElementById('generateBtn');
  btn.disabled = true;
  btn.innerHTML = '<span class="spinner"></span> Generating presentation...';
  showStatus('🤖 AI is creating your slide content... This takes 15-30 seconds.', 'loading');
  
  const formData = new FormData();
  formData.append('client_name', clientName);
  formData.append('company_name', companyName);
  formData.append('presentation_type', presType);
  formData.append('tone', tone);
  formData.append('font_style', fontStyle);
  formData.append('key_points', keyPoints);
  formData.append('num_slides', numSlides);
  
  const logoInput = document.getElementById('logoInput');
  if (logoInput.files.length) {
    formData.append('logo', logoInput.files[0]);
  }
  
  try {
    const res = await fetch('/api/generate', { method: 'POST', body: formData });
    const data = await res.json();
    
    if (data.success) {
      const downloadUrl = data.download_url + '?name=' + encodeURIComponent(data.filename);
      showStatus(`
        <div>✅ Presentation generated! ${data.slides_count} slides created.</div>
        <a href="${downloadUrl}" class="download-btn">📥 Download ${data.filename}</a>
      `, 'success');
    } else {
      showStatus('❌ ' + (data.error || 'Failed to generate'), 'error');
    }
  } catch(e) {
    showStatus('❌ Connection error: ' + e.message, 'error');
  }
  
  btn.disabled = false;
  btn.innerHTML = '⚡ Generate Presentation';
}

function showStatus(msg, type, targetId) {
  const el = document.getElementById(targetId || 'status');
  el.innerHTML = msg;
  el.className = 'status ' + type;
}

// ── Mode Switcher ──
const modes = ['create','polish','version','style','merge'];
const modeIds = {create:'modeCreate',polish:'modePolish',version:'modeVersion',style:'modeStyle',merge:'modeMerge'};
const tabIds = {create:'tabCreate',polish:'tabPolish',version:'tabVersion',style:'tabStyle',merge:'tabMerge'};

function switchMode(mode) {
  modes.forEach(m => {
    document.getElementById(modeIds[m]).style.display = m === mode ? 'block' : 'none';
    document.getElementById(tabIds[m]).style.background = m === mode ? 'linear-gradient(135deg,#6C5CE7,#5A4BD1)' : 'transparent';
    document.getElementById(tabIds[m]).style.color = m === mode ? '#fff' : '#8B8FA3';
  });
}

// ── Shared Helpers ──
function handleFileUpload(input, dropzoneId, textId) {
  if (!input.files.length) return;
  document.getElementById(textId).textContent = '✓ ' + input.files[0].name;
  document.getElementById(dropzoneId).classList.add('has-file');
}

function handleMultiFile(input) {
  if (!input.files.length) return;
  const names = Array.from(input.files).map(f => f.name).join(', ');
  document.getElementById('mergeText').textContent = '✓ ' + input.files.length + ' files: ' + names;
  document.getElementById('mergeDropzone').classList.add('has-file');
}

function makeDownloadLink(url, fname) {
  return '<a href="' + url + '?name=' + encodeURIComponent(fname) + '" class="download-btn">📥 ' + fname + '</a>';
}

// ── Polish Mode ──
const polishInstEl = document.getElementById('polishInstructions');
polishInstEl.addEventListener('input', () => {
  document.getElementById('polishCharCount').textContent = polishInstEl.value.length;
});

function handlePptx(input) {
  if (!input.files.length) return;
  document.getElementById('pptxText').textContent = '✓ ' + input.files[0].name;
  document.getElementById('pptxDropzone').classList.add('has-file');
}

function handlePolishLogo(input) {
  if (!input.files.length) return;
  document.getElementById('polishLogoText').textContent = '✓ ' + input.files[0].name;
}

function fillPolishExample(type) {
  const examples = {
    concise: "Make every slide more concise. Remove filler words. Keep bullets to 8-12 words max. Cut any redundant slides. Add a strong executive summary slide at the beginning.",
    visual: "Convert text-heavy slides into visual layouts. Add a stats slide with key metrics. Include a timeline for milestones. Use comparison tables instead of long paragraphs. Add an icon grid for features/services.",
    investor: "Restructure for investor pitch format: Problem → Solution → Market Size → Traction → Business Model → Team → Ask. Make numbers and metrics prominent. Add a competitive landscape comparison. End with a clear funding ask.",
    executive: "Condense to 8 slides max. Lead with the conclusion/recommendation. Use stats and charts over bullets. Remove technical details. Make every slide answerable in under 30 seconds. Add a decision-ready closing slide."
  };
  polishInstEl.value = examples[type] || '';
  document.getElementById('polishCharCount').textContent = polishInstEl.value.length;
}

async function polishDeck() {
  const pptxInput = document.getElementById('pptxInput');
  const instructions = polishInstEl.value.trim();
  if (!pptxInput.files.length) { showStatus('Please upload a PPTX file', 'error', 'polishStatus'); return; }
  if (!instructions) { showStatus('Please enter polishing instructions', 'error', 'polishStatus'); return; }

  const btn = document.getElementById('polishBtn');
  btn.disabled = true;
  btn.innerHTML = '<span class="spinner"></span> Polishing your deck...';
  showStatus('🤖 AI is reading your slides and applying changes... This takes 20-40 seconds.', 'loading', 'polishStatus');

  const formData = new FormData();
  formData.append('pptx_file', pptxInput.files[0]);
  formData.append('instructions', instructions);
  formData.append('tone', document.getElementById('polishTone').value);
  formData.append('font_style', document.getElementById('polishFontStyle').value);
  formData.append('num_slides', document.getElementById('polishNumSlides').value);
  formData.append('client_name', document.getElementById('polishClientName').value.trim());
  formData.append('company_name', document.getElementById('polishCompanyName').value.trim());
  const logoInput = document.getElementById('polishLogoInput');
  if (logoInput.files.length) formData.append('logo', logoInput.files[0]);

  try {
    const res = await fetch('/api/polish', { method: 'POST', body: formData });
    const data = await res.json();
    if (data.success) {
      showStatus('<div>✅ Deck polished! ' + data.original_slides + ' slides → ' + data.slides_count + ' slides.</div>' +
        makeDownloadLink(data.download_url, data.filename), 'success', 'polishStatus');
    } else {
      showStatus('❌ ' + (data.error || 'Failed'), 'error', 'polishStatus');
    }
  } catch(e) { showStatus('❌ ' + e.message, 'error', 'polishStatus'); }
  btn.disabled = false;
  btn.innerHTML = '✨ Polish Presentation';
}

// ── Audience Versions ──
async function generateVersions() {
  const input = document.getElementById('versionInput');
  if (!input.files.length) { showStatus('Please upload a PPTX file', 'error', 'versionStatus'); return; }

  const audiences = [];
  if (document.getElementById('vExec').checked) audiences.push('executive');
  if (document.getElementById('vDetailed').checked) audiences.push('detailed');
  if (document.getElementById('vInvestor').checked) audiences.push('investor');
  if (!audiences.length) { showStatus('Select at least one audience', 'error', 'versionStatus'); return; }

  const btn = document.getElementById('versionBtn');
  btn.disabled = true;
  btn.innerHTML = '<span class="spinner"></span> Generating ' + audiences.length + ' versions...';
  showStatus('🤖 Creating ' + audiences.length + ' audience versions... This takes 30-60 seconds.', 'loading', 'versionStatus');

  const formData = new FormData();
  formData.append('pptx_file', input.files[0]);
  audiences.forEach(a => formData.append('audiences', a));
  formData.append('client_name', document.getElementById('vClientName').value.trim());
  formData.append('company_name', document.getElementById('vCompanyName').value.trim());
  formData.append('font_style', document.getElementById('vFontStyle').value);
  const logo = document.getElementById('vLogoInput');
  if (logo.files.length) formData.append('logo', logo.files[0]);

  try {
    const res = await fetch('/api/version', { method: 'POST', body: formData });
    const data = await res.json();
    if (data.success) {
      let html = '<div>✅ ' + data.versions.length + ' versions created from ' + data.original_slides + ' original slides:</div>';
      data.versions.forEach(v => {
        if (v.error) {
          html += '<div style="margin-top:8px;color:var(--red)">❌ ' + v.audience + ': ' + v.error + '</div>';
        } else {
          html += '<div style="margin-top:8px">' + makeDownloadLink(v.download_url, v.filename) + ' <span style="color:var(--text2);font-size:12px">' + v.slides_count + ' slides</span></div>';
        }
      });
      showStatus(html, 'success', 'versionStatus');
    } else {
      showStatus('❌ ' + (data.error || 'Failed'), 'error', 'versionStatus');
    }
  } catch(e) { showStatus('❌ ' + e.message, 'error', 'versionStatus'); }
  btn.disabled = false;
  btn.innerHTML = '👥 Generate All Versions';
}

// ── Style Transfer ──
async function doStyleTransfer() {
  const contentInput = document.getElementById('contentInput');
  const refInput = document.getElementById('refInput');
  if (!contentInput.files.length) { showStatus('Please upload your content deck', 'error', 'styleStatus'); return; }
  if (!refInput.files.length) { showStatus('Please upload a reference/style deck', 'error', 'styleStatus'); return; }

  const btn = document.getElementById('styleBtn');
  btn.disabled = true;
  btn.innerHTML = '<span class="spinner"></span> Applying style...';
  showStatus('🤖 Analyzing reference style and restyling your content... 20-40 seconds.', 'loading', 'styleStatus');

  const formData = new FormData();
  formData.append('content_file', contentInput.files[0]);
  formData.append('reference_file', refInput.files[0]);
  formData.append('instructions', document.getElementById('styleInstructions').value.trim());
  formData.append('client_name', document.getElementById('sClientName').value.trim());
  formData.append('company_name', document.getElementById('sCompanyName').value.trim());
  formData.append('font_style', document.getElementById('sFontStyle').value);
  const logo = document.getElementById('sLogoInput');
  if (logo.files.length) formData.append('logo', logo.files[0]);

  try {
    const res = await fetch('/api/style-transfer', { method: 'POST', body: formData });
    const data = await res.json();
    if (data.success) {
      showStatus('<div>✅ Style applied! ' + data.content_slides + ' content slides restyled using ' + data.reference_slides + '-slide reference → ' + data.slides_count + ' slides.</div>' +
        makeDownloadLink(data.download_url, data.filename), 'success', 'styleStatus');
    } else {
      showStatus('❌ ' + (data.error || 'Failed'), 'error', 'styleStatus');
    }
  } catch(e) { showStatus('❌ ' + e.message, 'error', 'styleStatus'); }
  btn.disabled = false;
  btn.innerHTML = '🎭 Apply Style Transfer';
}

// ── Merge ──
async function doMerge() {
  const input = document.getElementById('mergeInput');
  if (!input.files.length || input.files.length < 2) { showStatus('Please upload at least 2 PPTX files', 'error', 'mergeStatus'); return; }

  const btn = document.getElementById('mergeBtn');
  btn.disabled = true;
  btn.innerHTML = '<span class="spinner"></span> Merging...';
  showStatus('🤖 Merging ' + input.files.length + ' decks... 20-40 seconds.', 'loading', 'mergeStatus');

  const formData = new FormData();
  Array.from(input.files).forEach(f => formData.append('pptx_files', f));
  formData.append('instructions', document.getElementById('mergeInstructions').value.trim());
  formData.append('client_name', document.getElementById('mClientName').value.trim());
  formData.append('company_name', document.getElementById('mCompanyName').value.trim());
  formData.append('font_style', document.getElementById('mFontStyle').value);
  formData.append('tone', document.getElementById('mTone').value);

  try {
    const res = await fetch('/api/merge', { method: 'POST', body: formData });
    const data = await res.json();
    if (data.success) {
      showStatus('<div>✅ Merged! ' + data.source_files + ' files (' + data.source_slides + ' slides) → ' + data.slides_count + ' cohesive slides.</div>' +
        makeDownloadLink(data.download_url, data.filename), 'success', 'mergeStatus');
    } else {
      showStatus('❌ ' + (data.error || 'Failed'), 'error', 'mergeStatus');
    }
  } catch(e) { showStatus('❌ ' + e.message, 'error', 'mergeStatus'); }
  btn.disabled = false;
  btn.innerHTML = '🔗 Merge Decks';
}

// ── Split ──
async function doSplit() {
  const input = document.getElementById('splitInput');
  if (!input.files.length) { showStatus('Please upload a PPTX file to split', 'error', 'mergeStatus'); return; }

  const btn = document.getElementById('splitBtn');
  btn.disabled = true;
  btn.innerHTML = '<span class="spinner"></span> Splitting...';
  showStatus('🤖 Creating Executive Summary + Full Detail versions... 30-50 seconds.', 'loading', 'mergeStatus');

  const formData = new FormData();
  formData.append('pptx_file', input.files[0]);
  formData.append('client_name', document.getElementById('mClientName').value.trim());
  formData.append('company_name', document.getElementById('mCompanyName').value.trim());
  formData.append('font_style', document.getElementById('mFontStyle').value);

  try {
    const res = await fetch('/api/split', { method: 'POST', body: formData });
    const data = await res.json();
    if (data.success) {
      let html = '<div>✅ Split! ' + data.original_slides + ' original slides into 2 versions:</div>';
      data.parts.forEach(p => {
        html += '<div style="margin-top:8px">' + makeDownloadLink(p.download_url, p.filename) + ' <span style="color:var(--text2);font-size:12px">' + p.part + ' · ' + p.slides_count + ' slides</span></div>';
      });
      showStatus(html, 'success', 'mergeStatus');
    } else {
      showStatus('❌ ' + (data.error || 'Failed'), 'error', 'mergeStatus');
    }
  } catch(e) { showStatus('❌ ' + e.message, 'error', 'mergeStatus'); }
  btn.disabled = false;
  btn.innerHTML = '✂️ Split Deck';
}
</script>
</body></html>"""

@app.route('/')
def index():
    return redirect('/welcome')

@app.route('/welcome')
def welcome():
    return render_template_string(LANDING_HTML)

@app.route('/login', methods=['GET'])
def login():
    if 'user_id' in session: return redirect('/create')
    return render_template('login.html')

@app.route('/register', methods=['GET'])
def register():
    if 'user_id' in session: return redirect('/create')
    return render_template('login.html', show_register=True)

@app.route('/logout')
def logout():
    session.clear()
    return redirect('/welcome')

# --- OTP API ---
@app.route('/api/auth/send-otp', methods=['POST'])
def send_otp():
    data = request.get_json()
    email = (data.get('email') or '').strip().lower()
    purpose = data.get('purpose', 'login')
    if not email or '@' not in email:
        return jsonify({"error": "Valid email required"}), 400
    conn = get_db()
    if not conn: return jsonify({"error": "Database not configured"}), 500
    cur = conn.cursor(cursor_factory=psycopg2.extras.RealDictCursor)
    cur.execute("""SELECT COUNT(*) as cnt FROM otp_codes
                   WHERE email=%s AND created_at > NOW() - INTERVAL '15 minutes'""", (email,))
    if cur.fetchone()['cnt'] >= 5:
        conn.close()
        return jsonify({"error": "Too many requests. Wait 15 minutes."}), 429
    if purpose == 'login':
        cur.execute('SELECT id FROM users WHERE email=%s', (email,))
        if not cur.fetchone():
            conn.close()
            return jsonify({"error": "No account found with this email"}), 404
    if purpose == 'register':
        cur.execute('SELECT id FROM users WHERE email=%s', (email,))
        if cur.fetchone():
            conn.close()
            return jsonify({"error": "Email already registered. Please sign in."}), 409
    cur.execute("UPDATE otp_codes SET used=TRUE WHERE email=%s AND purpose=%s AND used=FALSE", (email, purpose))
    code = generate_otp()
    expires = datetime.utcnow() + timedelta(minutes=5)
    cur.execute("INSERT INTO otp_codes (email, code, purpose, expires_at) VALUES (%s,%s,%s,%s)",
                (email, code, purpose, expires))
    conn.close()
    if send_otp_email(email, code, purpose):
        return jsonify({"success": True})
    return jsonify({"error": "Failed to send email"}), 500

@app.route('/api/auth/verify-otp', methods=['POST'])
def verify_otp():
    data = request.get_json()
    email = (data.get('email') or '').strip().lower()
    code = (data.get('code') or '').strip()
    purpose = data.get('purpose', 'login')
    if not email or not code or len(code) != 6:
        return jsonify({"error": "Email and 6-digit code required"}), 400
    conn = get_db()
    if not conn: return jsonify({"error": "Database not configured"}), 500
    cur = conn.cursor(cursor_factory=psycopg2.extras.RealDictCursor)
    cur.execute("""SELECT * FROM otp_codes
                   WHERE email=%s AND purpose=%s AND used=FALSE AND expires_at > NOW()
                   ORDER BY created_at DESC LIMIT 1""", (email, purpose))
    otp_rec = cur.fetchone()
    if not otp_rec:
        conn.close()
        return jsonify({"error": "Code expired. Request a new one."}), 400
    if otp_rec['attempts'] >= 3:
        cur.execute("UPDATE otp_codes SET used=TRUE WHERE id=%s", (otp_rec['id'],))
        conn.close()
        return jsonify({"error": "Too many attempts. Request a new code."}), 429
    cur.execute("UPDATE otp_codes SET attempts=attempts+1 WHERE id=%s", (otp_rec['id'],))
    if not secrets.compare_digest(code, otp_rec['code']):
        conn.close()
        remaining = 2 - otp_rec['attempts']
        return jsonify({"error": f"Invalid code. {remaining} attempt(s) remaining."}), 400
    cur.execute("UPDATE otp_codes SET used=TRUE WHERE id=%s", (otp_rec['id'],))
    if purpose == 'login':
        cur.execute('SELECT * FROM users WHERE email=%s', (email,))
        user = cur.fetchone()
        conn.close()
        if user:
            session['user_id'] = user['id']
            session['company_name'] = user.get('company_name', '')
            session.permanent = True
            return jsonify({"success": True, "redirect": "/create"})
        return jsonify({"error": "User not found"}), 404
    conn.close()
    return jsonify({"success": True, "verified": True})

@app.route('/api/auth/register', methods=['POST'])
def api_register():
    data = request.get_json()
    email = (data.get('email') or '').strip().lower()
    password = data.get('password', '')
    company = (data.get('company_name') or '').strip()
    currency = data.get('currency', 'USD')
    code = (data.get('code') or '').strip()
    if not email or not password or not company:
        return jsonify({"error": "All fields required"}), 400
    if len(password) < 8:
        return jsonify({"error": "Password must be at least 8 characters"}), 400
    if len(code) != 6:
        return jsonify({"error": "Valid 6-digit code required"}), 400
    conn = get_db()
    if not conn: return jsonify({"error": "Database not configured"}), 500
    cur = conn.cursor(cursor_factory=psycopg2.extras.RealDictCursor)
    cur.execute("""SELECT * FROM otp_codes
                   WHERE email=%s AND purpose='register' AND used=FALSE AND expires_at > NOW()
                   ORDER BY created_at DESC LIMIT 1""", (email,))
    otp_rec = cur.fetchone()
    if not otp_rec or not secrets.compare_digest(code, otp_rec['code']):
        conn.close()
        return jsonify({"error": "Invalid or expired code"}), 400
    if otp_rec['attempts'] >= 3:
        conn.close()
        return jsonify({"error": "Too many attempts. Request a new code."}), 429
    cur.execute("UPDATE otp_codes SET used=TRUE WHERE id=%s", (otp_rec['id'],))
    cur.execute('SELECT id FROM users WHERE email=%s', (email,))
    if cur.fetchone():
        conn.close()
        return jsonify({"error": "Email already registered"}), 409
    cur.execute('SELECT COUNT(*) as cnt FROM users')
    is_first = cur.fetchone()['cnt'] == 0
    try:
        cur.execute('INSERT INTO users (email, password_hash, company_name, currency, is_superadmin) VALUES (%s,%s,%s,%s,%s) RETURNING id',
                    (email, hash_pw(password), company, currency, is_first))
        user_id = cur.fetchone()['id']
        conn.close()
        session['user_id'] = user_id
        session['company_name'] = company
        session.permanent = True
        register_with_hub(company, email, currency)
        return jsonify({"success": True, "redirect": "/create"})
    except psycopg2.errors.UniqueViolation:
        conn.close()
        return jsonify({"error": "Email already registered"}), 409

@app.route('/create')
@login_required
def create():
    is_admin = False
    try:
        conn = get_db()
        if conn:
            cur = conn.cursor(cursor_factory=psycopg2.extras.RealDictCursor)
            cur.execute('SELECT is_superadmin FROM users WHERE id=%s', (session.get('user_id'),))
            u = cur.fetchone()
            is_admin = u.get('is_superadmin', False) if u else False
            conn.close()
    except: pass
    return render_template_string(MAIN_HTML, is_admin=is_admin)

@app.route('/admin')
@login_required
def admin_dashboard():
    user_id = session.get('user_id')
    try:
        conn = get_db()
        if not conn: return 'Database not configured', 500
        cur = conn.cursor(cursor_factory=psycopg2.extras.RealDictCursor)

        # Check superadmin
        cur.execute('SELECT * FROM users WHERE id=%s', (user_id,))
        user = cur.fetchone()
        if not user or not user.get('is_superadmin'):
            flash('Access denied', 'error')
            return redirect('/create')

        # Platform stats
        cur.execute('''SELECT
            (SELECT COUNT(*) FROM users) as total_users,
            (SELECT COUNT(*) FROM usage_log) as total_presentations,
            (SELECT COALESCE(SUM(slides),0) FROM usage_log) as total_slides,
            (SELECT COUNT(*) FROM users WHERE created_at > CURRENT_DATE - INTERVAL '7 days') as new_users_7d,
            (SELECT COUNT(*) FROM usage_log WHERE created_at > CURRENT_DATE - INTERVAL '7 days') as presentations_7d
        ''')
        platform = cur.fetchone()

        # All users with usage counts
        cur.execute('''SELECT u.id, u.email, u.company_name, u.currency, u.is_superadmin, u.created_at,
                      COUNT(l.id) as presentation_count,
                      COALESCE(SUM(l.slides),0) as total_slides,
                      MAX(l.created_at) as last_activity
                      FROM users u LEFT JOIN usage_log l ON u.id = l.user_id
                      GROUP BY u.id ORDER BY u.created_at DESC''')
        users = cur.fetchall()

        # Recent activity (last 50)
        cur.execute('''SELECT l.*, u.email, u.company_name
                      FROM usage_log l JOIN users u ON l.user_id = u.id
                      ORDER BY l.created_at DESC LIMIT 50''')
        activity = cur.fetchall()

        conn.close()
        return render_template_string(ADMIN_HTML, user=user, platform=platform,
                                     users=users, activity=activity)
    except Exception as e:
        return f'Error: {e}', 500

ADMIN_HTML = """<!DOCTYPE html>
<html lang="en"><head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>ProposalSnap — Admin Dashboard</title>
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&display=swap" rel="stylesheet">
<style>
:root{--bg:#0B0F1A;--surface:#131829;--border:rgba(255,255,255,0.08);--text:#F0F0F5;--text2:#8B8FA3;--accent:#6C5CE7;--accent2:#A78BFA;--green:#00D2A0;--red:#FF6B6B;--orange:#f59e0b}
*{margin:0;padding:0;box-sizing:border-box}
body{font-family:'Inter',sans-serif;background:var(--bg);color:var(--text);min-height:100vh}
.wrap{max-width:1100px;margin:0 auto;padding:24px}
.hdr{display:flex;align-items:center;justify-content:space-between;margin-bottom:32px;flex-wrap:wrap;gap:12px}
.hdr h1{font-size:24px;font-weight:800;color:#fff}
.hdr h1 span{color:var(--accent)}
.hdr-links a{font-size:13px;color:var(--text2);text-decoration:none;padding:8px 16px;border:1px solid var(--border);border-radius:8px;font-weight:600;margin-left:8px}
.hdr-links a:hover{border-color:var(--accent);color:var(--accent)}
.stats{display:grid;grid-template-columns:repeat(auto-fit,minmax(180px,1fr));gap:14px;margin-bottom:32px}
.stat{background:var(--surface);border:1px solid var(--border);border-radius:14px;padding:20px}
.stat .label{font-size:11px;font-weight:700;text-transform:uppercase;letter-spacing:1px;color:var(--text2);margin-bottom:6px}
.stat .value{font-size:28px;font-weight:800;color:#fff}
.stat .sub{font-size:12px;color:var(--text2);margin-top:4px}
.section{margin-bottom:32px}
.section h2{font-size:16px;font-weight:700;color:#fff;margin-bottom:14px}
.tbl{width:100%;border-collapse:collapse;background:var(--surface);border:1px solid var(--border);border-radius:14px;overflow:hidden}
.tbl th{font-size:11px;font-weight:700;text-transform:uppercase;letter-spacing:.5px;color:var(--text2);padding:12px 16px;text-align:left;border-bottom:1px solid var(--border);background:rgba(255,255,255,.02)}
.tbl td{padding:10px 16px;font-size:13px;color:var(--text);border-bottom:1px solid rgba(255,255,255,.04)}
.tbl tr:last-child td{border-bottom:none}
.badge{display:inline-block;padding:2px 8px;border-radius:6px;font-size:11px;font-weight:700}
.badge-admin{background:rgba(167,139,250,.15);color:var(--accent2)}
.badge-new{background:rgba(0,210,160,.15);color:var(--green)}
.empty{text-align:center;padding:40px;color:var(--text2);font-size:14px}
@media(max-width:600px){.stats{grid-template-columns:1fr 1fr}.tbl{font-size:12px}.tbl th,.tbl td{padding:8px 10px}}
</style></head>
<body>
<div class="wrap">
<div class="hdr">
<h1>Proposal<span>Snap</span> Admin</h1>
<div class="hdr-links">
<a href="/create">← Back to Tool</a>
<a href="/logout" style="color:var(--red);border-color:rgba(255,107,107,.3)">Sign Out</a>
</div>
</div>

<div class="stats">
<div class="stat"><div class="label">Total Users</div><div class="value">{{ platform.total_users }}</div><div class="sub">+{{ platform.new_users_7d }} this week</div></div>
<div class="stat"><div class="label">Presentations</div><div class="value">{{ platform.total_presentations }}</div><div class="sub">{{ platform.presentations_7d }} this week</div></div>
<div class="stat"><div class="label">Total Slides</div><div class="value">{{ platform.total_slides }}</div></div>
<div class="stat"><div class="label">Avg per User</div><div class="value">{{ '%.1f' | format(platform.total_presentations / platform.total_users if platform.total_users > 0 else 0) }}</div><div class="sub">presentations / user</div></div>
</div>

<div class="section">
<h2>Registered Users ({{ users|length }})</h2>
{% if users %}
<table class="tbl">
<tr><th>Email</th><th>Company</th><th>Presentations</th><th>Slides</th><th>Last Active</th><th>Joined</th></tr>
{% for u in users %}
<tr>
<td>{{ u.email }} {% if u.is_superadmin %}<span class="badge badge-admin">Admin</span>{% endif %}</td>
<td>{{ u.company_name or '—' }}</td>
<td>{{ u.presentation_count }}</td>
<td>{{ u.total_slides }}</td>
<td>{{ u.last_activity.strftime('%d %b %Y') if u.last_activity else '—' }}</td>
<td>{{ u.created_at.strftime('%d %b %Y') }}</td>
</tr>
{% endfor %}
</table>
{% else %}
<div class="empty">No users yet</div>
{% endif %}
</div>

<div class="section">
<h2>Recent Activity</h2>
{% if activity %}
<table class="tbl">
<tr><th>User</th><th>Company</th><th>Presentation</th><th>Slides</th><th>When</th></tr>
{% for a in activity %}
<tr>
<td>{{ a.email }}</td>
<td>{{ a.company_name or '—' }}</td>
<td>{{ a.title or '—' }}</td>
<td>{{ a.slides }}</td>
<td>{{ a.created_at.strftime('%d %b %Y %H:%M') }}</td>
</tr>
{% endfor %}
</table>
{% else %}
<div class="empty">No presentations generated yet</div>
{% endif %}
</div>
</div>
</body></html>"""

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    print(f"\n{'='*50}")
    print(f"  🎯 ProposalSnap is running!")
    print(f"{'='*50}")
    print(f"  Open: http://localhost:{port}")
    print(f"{'='*50}\n")
    app.run(host='0.0.0.0', port=port, debug=True)

@app.route('/demo')
def demo():
    hub_url = os.environ.get('HUB_URL', 'https://snapsuite.up.railway.app')
    return render_template_string(DEMO_GALLERY_HTML, hub_url=hub_url)

DEMO_GALLERY_HTML = r'''<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>ProposalSnap — Demo Gallery</title>
<link href="https://fonts.googleapis.com/css2?family=DM+Sans:wght@400;500;600;700;800;900&family=Playfair+Display:wght@600;700;800;900&family=JetBrains+Mono:wght@400;500;600&display=swap" rel="stylesheet">
<style>
:root{--bg:#06080F;--surface:#0D1117;--card:#131920;--border:#1C2433;--border2:#2A3548;--text:#E2E8F0;--text2:#6B7A90;
--blue:#3B82F6;--green:#4ADE80;--red:#F87171;--orange:#F59E0B;--purple:#A78BFA;--teal:#2DD4BF;--pink:#F472B6;--cyan:#22D3EE;--indigo:#818CF8}
*{margin:0;padding:0;box-sizing:border-box}
body{font-family:'DM Sans',sans-serif;background:var(--bg);color:var(--text);min-height:100vh}
body::before{content:'';position:fixed;inset:0;background:radial-gradient(ellipse at 20% 0%,rgba(59,130,246,.06) 0%,transparent 60%),radial-gradient(ellipse at 80% 100%,rgba(167,139,250,.05) 0%,transparent 50%);pointer-events:none;z-index:0}

/* Topbar */
.topbar{position:sticky;top:0;z-index:100;background:rgba(6,8,15,.85);backdrop-filter:blur(24px);border-bottom:1px solid var(--border);padding:14px 32px;display:flex;align-items:center;justify-content:space-between}
.topbar h1{font-size:22px;font-weight:900;color:#fff;letter-spacing:-.5px}
.topbar h1 span{background:linear-gradient(135deg,var(--blue),var(--purple));-webkit-background-clip:text;-webkit-text-fill-color:transparent}
.topbar-right{display:flex;gap:8px}
.topbar-right a{font-size:13px;color:var(--text2);text-decoration:none;padding:8px 16px;border-radius:8px;font-weight:600;transition:.2s}
.topbar-right a:hover{color:#fff;background:rgba(255,255,255,.06)}
.topbar-right .cta{background:linear-gradient(135deg,var(--blue),var(--purple));color:#fff;font-weight:700}

/* Hero */
.hero-section{padding:48px 32px 16px;max-width:1200px;margin:0 auto;position:relative;z-index:1}
.hero-section h2{font-size:36px;font-weight:900;color:#fff;letter-spacing:-1.5px;margin-bottom:8px}
.hero-section .sub{font-size:16px;color:var(--text2);line-height:1.6;max-width:600px}
.hero-section .badge{display:inline-flex;align-items:center;gap:6px;padding:6px 14px;background:rgba(59,130,246,.08);border:1px solid rgba(59,130,246,.15);border-radius:50px;font-size:12px;font-weight:700;color:var(--blue);margin-bottom:16px}
.hero-section .badge .dot{width:6px;height:6px;border-radius:50%;background:var(--green);animation:pulse 2s infinite}
@keyframes pulse{0%,100%{opacity:1}50%{opacity:.3}}

/* Main */
.main{max-width:1200px;margin:0 auto;padding:32px;position:relative;z-index:1}

/* Proposal section */
.proposal{margin-bottom:48px;position:relative}
.proposal-label{font-family:'JetBrains Mono',monospace;font-size:11px;text-transform:uppercase;letter-spacing:2px;font-weight:600;margin-bottom:12px}

/* Proposal header card */
.proposal-head{background:var(--card);border:1px solid var(--border);border-radius:16px;padding:28px 32px;margin-bottom:16px;display:flex;justify-content:space-between;align-items:center;position:relative;overflow:hidden}
.proposal-head::before{content:'';position:absolute;top:0;left:0;right:0;height:3px}
.proposal-head.theme-corporate::before{background:linear-gradient(90deg,#e94560,#c62a88,#0f3460)}
.proposal-head.theme-wedding::before{background:linear-gradient(90deg,#f5c6a5,#c62a88,#801336)}
.proposal-head.theme-workshop::before{background:linear-gradient(90deg,#64ffda,#3B82F6,#0a192f)}
.ph-left h3{font-size:22px;font-weight:800;color:#fff;letter-spacing:-.5px;margin-bottom:6px}
.ph-left .ph-meta{font-size:13px;color:var(--text2);display:flex;align-items:center;gap:12px;flex-wrap:wrap}
.ph-left .ph-meta .sep{color:var(--border2)}
.ph-right{text-align:right}
.ph-right .price{font-size:28px;font-weight:900;color:#fff;letter-spacing:-1px;margin-bottom:4px}
.status{display:inline-flex;align-items:center;gap:4px;padding:4px 12px;border-radius:20px;font-size:11px;font-weight:700;letter-spacing:.3px}
.status.sent{background:rgba(59,130,246,.12);color:var(--blue);border:1px solid rgba(59,130,246,.2)}
.status.accepted{background:rgba(74,222,128,.1);color:var(--green);border:1px solid rgba(74,222,128,.2)}
.status.draft{background:rgba(107,122,144,.1);color:var(--text2);border:1px solid rgba(107,122,144,.15)}
.status::before{content:'';width:6px;height:6px;border-radius:50%;background:currentColor}

/* Slides scroll */
.slides-wrapper{position:relative}
.slides-scroll{display:flex;gap:14px;overflow-x:auto;padding:8px 4px 20px;scroll-snap-type:x mandatory;-webkit-overflow-scrolling:touch;scrollbar-width:none}
.slides-scroll::-webkit-scrollbar{display:none}
.slides-scroll::after{content:'';min-width:20px;flex-shrink:0}

/* Individual slide */
.slide{min-width:300px;max-width:300px;aspect-ratio:16/9;border-radius:12px;flex-shrink:0;scroll-snap-align:start;position:relative;overflow:hidden;cursor:default;transition:transform .3s,box-shadow .3s;border:1px solid rgba(255,255,255,.06)}
.slide:hover{transform:translateY(-4px) scale(1.02);box-shadow:0 16px 48px rgba(0,0,0,.5)}
.slide-inner{position:absolute;inset:0;padding:22px 24px;display:flex;flex-direction:column;z-index:2}
.slide-num{position:absolute;top:10px;right:12px;font-family:'JetBrains Mono',monospace;font-size:10px;font-weight:600;opacity:.35;z-index:3}

/* Slide decorative elements */
.slide-deco{position:absolute;inset:0;z-index:1;overflow:hidden;pointer-events:none}
.slide-deco .circle{position:absolute;border-radius:50%;opacity:.08}
.slide-deco .line{position:absolute;height:1px;opacity:.1}
.slide-deco .corner{position:absolute;width:40px;height:40px;opacity:.12}
.slide-deco .corner::before,.slide-deco .corner::after{content:'';position:absolute;background:currentColor}
.slide-deco .corner.tl{top:12px;left:12px}.slide-deco .corner.tl::before{top:0;left:0;width:16px;height:1.5px}.slide-deco .corner.tl::after{top:0;left:0;width:1.5px;height:16px}
.slide-deco .corner.br{bottom:12px;right:12px}.slide-deco .corner.br::before{bottom:0;right:0;width:16px;height:1.5px}.slide-deco .corner.br::after{bottom:0;right:0;width:1.5px;height:16px}
.slide-deco .dots{position:absolute;display:grid;grid-template-columns:repeat(5,4px);gap:8px;opacity:.08}
.slide-deco .dots span{width:4px;height:4px;border-radius:50%;background:currentColor}

/* Slide title styles */
.slide h4{font-family:'Playfair Display',serif;font-weight:800;line-height:1.2;margin-bottom:8px;position:relative}
.slide .slide-label{font-family:'JetBrains Mono',monospace;font-size:9px;text-transform:uppercase;letter-spacing:2px;opacity:.5;margin-bottom:auto}
.slide p{font-size:11px;line-height:1.65;opacity:.8}
.slide .tag{font-family:'JetBrains Mono',monospace;font-size:8px;padding:3px 8px;border-radius:4px;letter-spacing:1.5px;text-transform:uppercase;margin-top:auto;display:inline-block;width:fit-content}
.slide .divider{width:32px;height:2px;margin:8px 0;border-radius:1px}
.slide .icon-row{display:flex;gap:6px;margin-top:8px}
.slide .icon-row span{width:24px;height:24px;border-radius:6px;display:flex;align-items:center;justify-content:center;font-size:11px;background:rgba(255,255,255,.1)}

/* Theme: Corporate Navy (Brand Identity) */
.theme-corp{background:linear-gradient(145deg,#0a1628,#162240)}
.theme-corp-alt{background:linear-gradient(145deg,#162240,#0f3460)}
.theme-corp-accent{background:linear-gradient(145deg,#0f3460,#1a1a40)}
.theme-corp-dark{background:linear-gradient(145deg,#0a0f1f,#0a1628)}
.theme-corp h4{color:#e94560;font-size:16px}
.theme-corp p,.theme-corp-alt p,.theme-corp-accent p,.theme-corp-dark p{color:#b8c5d6}
.theme-corp .tag{background:rgba(233,69,96,.15);color:#e94560;border:1px solid rgba(233,69,96,.2)}
.theme-corp .divider,.theme-corp-alt .divider,.theme-corp-accent .divider{background:#e94560}
.theme-corp-alt h4,.theme-corp-accent h4,.theme-corp-dark h4{color:#e94560;font-size:16px}
.theme-corp .slide-deco .circle,.theme-corp-alt .slide-deco .circle{background:#e94560}
.theme-corp .slide-deco .line,.theme-corp-alt .slide-deco .line{background:#e94560}
.theme-corp .slide-deco .corner,.theme-corp-alt .slide-deco .corner,.theme-corp-accent .slide-deco .corner{color:#e94560}
.theme-corp .slide-deco .dots span,.theme-corp-alt .slide-deco .dots span{background:#e94560}

/* Theme: Wedding (Burgundy/Gold) */
.theme-wed{background:linear-gradient(145deg,#1f0a1a,#3d1132)}
.theme-wed-alt{background:linear-gradient(145deg,#3d1132,#5c1a4a)}
.theme-wed-accent{background:linear-gradient(145deg,#2a0e24,#1f0a1a)}
.theme-wed h4,.theme-wed-alt h4,.theme-wed-accent h4{color:#f5c6a5;font-size:16px}
.theme-wed p,.theme-wed-alt p,.theme-wed-accent p{color:#d4a9b8}
.theme-wed .tag{background:rgba(245,198,165,.1);color:#f5c6a5;border:1px solid rgba(245,198,165,.2)}
.theme-wed .divider,.theme-wed-alt .divider{background:linear-gradient(90deg,#f5c6a5,#c62a88)}
.theme-wed .slide-deco .circle,.theme-wed-alt .slide-deco .circle{background:#c62a88}
.theme-wed .slide-deco .corner,.theme-wed-alt .slide-deco .corner,.theme-wed-accent .slide-deco .corner{color:#f5c6a5}
.theme-wed .slide-deco .dots span,.theme-wed-alt .slide-deco .dots span{background:#f5c6a5}

/* Theme: Workshop (Teal/Mint) */
.theme-wk{background:linear-gradient(145deg,#06111f,#0a192f)}
.theme-wk-alt{background:linear-gradient(145deg,#0a192f,#112240)}
.theme-wk-accent{background:linear-gradient(145deg,#112240,#0a192f)}
.theme-wk h4,.theme-wk-alt h4,.theme-wk-accent h4{color:#64ffda;font-size:16px}
.theme-wk p,.theme-wk-alt p,.theme-wk-accent p{color:#8892b0}
.theme-wk .tag{background:rgba(100,255,218,.08);color:#64ffda;border:1px solid rgba(100,255,218,.15)}
.theme-wk .divider,.theme-wk-alt .divider{background:#64ffda}
.theme-wk .slide-deco .circle,.theme-wk-alt .slide-deco .circle{background:#64ffda}
.theme-wk .slide-deco .corner,.theme-wk-alt .slide-deco .corner,.theme-wk-accent .slide-deco .corner{color:#64ffda}
.theme-wk .slide-deco .dots span,.theme-wk-alt .slide-deco .dots span{background:#64ffda}

/* Title slide special */
.slide-title{justify-content:center;align-items:center;text-align:center}
.slide-title h4{font-size:20px!important;margin-bottom:10px}
.slide-title .company-from{font-family:'JetBrains Mono',monospace;font-size:10px;letter-spacing:1px;opacity:.5;margin-top:8px}

/* End slide special */
.slide-end{justify-content:center;align-items:center;text-align:center}
.slide-end h4{font-size:18px!important}
.slide-end .contact{font-family:'JetBrains Mono',monospace;font-size:10px;letter-spacing:.5px;opacity:.6;margin-top:10px}

/* Slide with metrics */
.metric-row{display:flex;gap:12px;margin-top:10px}
.metric{flex:1;text-align:center;padding:8px 4px;background:rgba(255,255,255,.04);border-radius:6px;border:1px solid rgba(255,255,255,.06)}
.metric .mv{font-size:16px;font-weight:800;color:#fff}
.metric .ml{font-size:8px;text-transform:uppercase;letter-spacing:1px;opacity:.5;margin-top:2px}

/* Timeline visual */
.timeline-row{display:flex;align-items:center;gap:0;margin-top:10px}
.tl-step{flex:1;text-align:center;position:relative;padding-top:12px}
.tl-step::before{content:'';position:absolute;top:4px;left:50%;width:8px;height:8px;border-radius:50%;transform:translateX(-50%)}
.tl-step::after{content:'';position:absolute;top:7px;left:50%;width:100%;height:1.5px;opacity:.2}
.tl-step:last-child::after{display:none}
.tl-step span{font-size:8px;line-height:1.3;display:block;opacity:.6}

/* Bullet list in slides */
.slide-list{list-style:none;padding:0;margin:6px 0}
.slide-list li{font-size:10px;padding:3px 0 3px 14px;position:relative;opacity:.75;line-height:1.5}
.slide-list li::before{content:'';position:absolute;left:0;top:8px;width:5px;height:5px;border-radius:50%}

/* Proposal footer */
.proposal-foot{display:flex;justify-content:space-between;align-items:center;padding:14px 0;margin-top:4px}
.proposal-foot .pf-tags{display:flex;gap:8px}
.proposal-foot .pf-tag{font-family:'JetBrains Mono',monospace;font-size:10px;color:var(--text2);background:var(--card);border:1px solid var(--border);padding:5px 12px;border-radius:6px;display:flex;align-items:center;gap:5px}
.proposal-foot a{font-size:13px;color:var(--blue);text-decoration:none;font-weight:700;transition:.2s}
.proposal-foot a:hover{color:#fff}

/* Scroll hint */
.scroll-hint{display:flex;align-items:center;gap:6px;font-size:11px;color:var(--text2);margin-bottom:8px}
.scroll-hint .arrow{animation:nudge 2s infinite;display:inline-block}
@keyframes nudge{0%,100%{transform:translateX(0)}50%{transform:translateX(6px)}}

/* Responsive */
@media(max-width:768px){
    .topbar{padding:12px 16px}
    .hero-section,.main{padding-left:16px;padding-right:16px}
    .proposal-head{flex-direction:column;gap:16px;text-align:left}
    .ph-right{text-align:left}
    .slide{min-width:260px;max-width:260px}
}
</style>
</head>
<body>

<div class="topbar">
    <a href="/" style="text-decoration:none;color:inherit"><h1>Proposal<span>Snap</span></h1></a>
    <div class="topbar-right">
        <a href="{{ hub_url }}">← SnapSuite</a>
        <a href="/" class="cta">+ Create New Proposal</a>
    </div>
</div>

<div class="hero-section">
    <div class="badge"><span class="dot"></span> AI-Generated Proposals</div>
    <h2>Sample Proposal Gallery</h2>
    <p class="sub">Each of these professional pitch decks was generated by AI in under 30 seconds. Upload your logo, describe your project, and ProposalSnap handles the rest.</p>
</div>

<div class="main">

<!-- ═══ PROPOSAL 1: Brand Identity ═══ -->
<div class="proposal">
    <div class="proposal-label" style="color:#e94560">● PROPOSAL 01</div>
    <div class="proposal-head theme-corporate">
        <div class="ph-left">
            <h3>Brand Identity Package</h3>
            <div class="ph-meta">
                <span>Bloom Studio → Varnam Artboutique</span>
                <span class="sep">·</span>
                <span>Feb 10, 2026</span>
                <span class="sep">·</span>
                <span>12 slides</span>
            </div>
        </div>
        <div class="ph-right">
            <div class="price">₹1,85,000</div>
            <span class="status sent">Sent</span>
        </div>
    </div>
    <div class="scroll-hint"><span>Scroll slides</span> <span class="arrow">→</span></div>
    <div class="slides-wrapper">
        <div class="slides-scroll">
            <!-- Slide 1: Title -->
            <div class="slide theme-corp">
                <div class="slide-deco">
                    <div class="corner tl" style="color:#e94560"></div>
                    <div class="corner br" style="color:#e94560"></div>
                    <div class="circle" style="width:200px;height:200px;right:-60px;bottom:-60px;background:#e94560"></div>
                    <div class="line" style="width:60%;bottom:30%;left:20%;background:linear-gradient(90deg,transparent,#e94560,transparent)"></div>
                </div>
                <div class="slide-inner slide-title">
                    <span class="tag">CORPORATE PROPOSAL</span>
                    <h4 style="margin-top:12px">Brand Identity<br>Package</h4>
                    <div class="divider" style="margin:8px auto"></div>
                    <p style="font-size:12px;color:#eee">Prepared for <strong>Varnam Artboutique</strong></p>
                    <span class="company-from">BY BLOOM STUDIO</span>
                </div>
                <div class="slide-num">01</div>
            </div>

            <!-- Slide 2: The Challenge -->
            <div class="slide theme-corp-alt">
                <div class="slide-deco">
                    <div class="corner tl" style="color:#e94560"></div>
                    <div class="dots" style="right:16px;bottom:16px;color:#e94560"><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span></div>
                </div>
                <div class="slide-inner">
                    <span class="slide-label">THE CHALLENGE</span>
                    <h4>Why Rebrand Now?</h4>
                    <div class="divider"></div>
                    <p>Varnam Artboutique is expanding into international markets. Current branding doesn't reflect their premium positioning or rich cultural heritage. Competitors are investing heavily in design.</p>
                    <div class="metric-row">
                        <div class="metric"><div class="mv" style="color:#e94560">73%</div><div class="ml">Need Update</div></div>
                        <div class="metric"><div class="mv" style="color:#e94560">2.4×</div><div class="ml">Brand Recall</div></div>
                    </div>
                </div>
                <div class="slide-num">02</div>
            </div>

            <!-- Slide 3: Our Approach -->
            <div class="slide theme-corp-accent">
                <div class="slide-deco">
                    <div class="corner tl" style="color:#e94560"></div>
                    <div class="corner br" style="color:#e94560"></div>
                    <div class="circle" style="width:120px;height:120px;left:-30px;top:-30px;background:#e94560"></div>
                </div>
                <div class="slide-inner">
                    <span class="slide-label">OUR APPROACH</span>
                    <h4>6-Phase Process</h4>
                    <div class="divider"></div>
                    <ul class="slide-list">
                        <li style="--c:#e94560">Discovery &amp; Brand Audit</li>
                        <li style="--c:#e94560">Concept Development</li>
                        <li style="--c:#e94560">Visual Design System</li>
                        <li style="--c:#e94560">Brand Guidelines</li>
                        <li style="--c:#e94560">Collateral Design</li>
                        <li style="--c:#e94560">Handoff &amp; Training</li>
                    </ul>
                </div>
                <div class="slide-num">03</div>
            </div>

            <!-- Slide 4: Logo -->
            <div class="slide theme-corp">
                <div class="slide-deco">
                    <div class="dots" style="left:16px;top:16px;color:#e94560"><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span></div>
                    <div class="line" style="width:40%;top:50%;right:0;background:#e94560"></div>
                </div>
                <div class="slide-inner">
                    <span class="slide-label">DELIVERABLE</span>
                    <h4>Logo Redesign</h4>
                    <div class="divider"></div>
                    <p>Modern mark blending Madhubani art motifs with contemporary typography. Primary, secondary, and icon variations.</p>
                    <div class="icon-row">
                        <span>🎨</span><span>✏️</span><span>📐</span>
                    </div>
                </div>
                <div class="slide-num">04</div>
            </div>

            <!-- Slide 5: Color Palette -->
            <div class="slide theme-corp-alt">
                <div class="slide-deco">
                    <div class="corner tl" style="color:#e94560"></div>
                    <div class="circle" style="width:80px;height:80px;right:20px;top:20px;background:#c62a88;opacity:.12"></div>
                    <div class="circle" style="width:60px;height:60px;right:60px;top:50px;background:#e94560;opacity:.1"></div>
                </div>
                <div class="slide-inner">
                    <span class="slide-label">VISUAL SYSTEM</span>
                    <h4>Color Palette</h4>
                    <div class="divider"></div>
                    <p>Rich earth tones paired with vibrant accents inspired by traditional Indian textiles and natural dyes.</p>
                    <div style="display:flex;gap:6px;margin-top:10px">
                        <div style="width:28px;height:28px;border-radius:6px;background:#e94560;border:2px solid rgba(255,255,255,.1)"></div>
                        <div style="width:28px;height:28px;border-radius:6px;background:#c62a88;border:2px solid rgba(255,255,255,.1)"></div>
                        <div style="width:28px;height:28px;border-radius:6px;background:#0f3460;border:2px solid rgba(255,255,255,.1)"></div>
                        <div style="width:28px;height:28px;border-radius:6px;background:#f5c6a5;border:2px solid rgba(255,255,255,.1)"></div>
                        <div style="width:28px;height:28px;border-radius:6px;background:#1a1a2e;border:2px solid rgba(255,255,255,.1)"></div>
                    </div>
                </div>
                <div class="slide-num">05</div>
            </div>

            <!-- Slide 6: Typography -->
            <div class="slide theme-corp-accent">
                <div class="slide-deco">
                    <div class="corner br" style="color:#e94560"></div>
                </div>
                <div class="slide-inner">
                    <span class="slide-label">VISUAL SYSTEM</span>
                    <h4>Typography</h4>
                    <div class="divider"></div>
                    <div style="margin-top:4px">
                        <div style="font-family:'Playfair Display',serif;font-size:18px;color:#fff;font-weight:700">Aa Heading</div>
                        <div style="font-size:10px;opacity:.5;margin-bottom:6px">Playfair Display · Serif</div>
                        <div style="font-family:'DM Sans',sans-serif;font-size:12px;color:#ccc">Aa Body text sample</div>
                        <div style="font-size:10px;opacity:.5;margin-bottom:6px">DM Sans · Sans-serif</div>
                        <div style="font-size:12px;color:#ccc">आ देवनागरी</div>
                        <div style="font-size:10px;opacity:.5">Matching Hindi family</div>
                    </div>
                </div>
                <div class="slide-num">06</div>
            </div>

            <!-- Slide 7: Collateral -->
            <div class="slide theme-corp">
                <div class="slide-deco">
                    <div class="dots" style="right:16px;top:16px;color:#e94560"><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span></div>
                    <div class="corner tl" style="color:#e94560"></div>
                </div>
                <div class="slide-inner">
                    <span class="slide-label">DELIVERABLES</span>
                    <h4>Collateral Design</h4>
                    <div class="divider"></div>
                    <ul class="slide-list">
                        <li style="--c:#e94560">Business cards &amp; letterhead</li>
                        <li style="--c:#e94560">Packaging inserts &amp; labels</li>
                        <li style="--c:#e94560">Exhibition banners (3 sizes)</li>
                        <li style="--c:#e94560">Social media templates</li>
                        <li style="--c:#e94560">WhatsApp catalog design</li>
                    </ul>
                </div>
                <div class="slide-num">07</div>
            </div>

            <!-- Slide 8: Digital -->
            <div class="slide theme-corp-alt">
                <div class="slide-deco">
                    <div class="corner tl" style="color:#e94560"></div>
                    <div class="corner br" style="color:#e94560"></div>
                </div>
                <div class="slide-inner">
                    <span class="slide-label">DIGITAL</span>
                    <h4>Online Presence</h4>
                    <div class="divider"></div>
                    <p>Website UI kit, email templates, Instagram grid layout, Google Business profile assets, and e-commerce product page templates.</p>
                </div>
                <div class="slide-num">08</div>
            </div>

            <!-- Slide 9: Guidelines -->
            <div class="slide theme-corp-accent">
                <div class="slide-deco">
                    <div class="circle" style="width:160px;height:160px;right:-40px;bottom:-40px;background:#e94560"></div>
                </div>
                <div class="slide-inner">
                    <span class="slide-label">DELIVERABLE</span>
                    <h4>Brand Guidelines</h4>
                    <div class="divider"></div>
                    <p>60-page comprehensive guide: logo usage, do's &amp; don'ts, spacing rules, color specs (CMYK, RGB, Pantone), voice &amp; tone.</p>
                    <div class="metric-row">
                        <div class="metric"><div class="mv" style="color:#e94560">60</div><div class="ml">Pages</div></div>
                        <div class="metric"><div class="mv" style="color:#e94560">3</div><div class="ml">Formats</div></div>
                    </div>
                </div>
                <div class="slide-num">09</div>
            </div>

            <!-- Slide 10: Timeline -->
            <div class="slide theme-corp">
                <div class="slide-deco">
                    <div class="corner tl" style="color:#e94560"></div>
                    <div class="line" style="width:80%;bottom:40px;left:10%;background:#e94560"></div>
                </div>
                <div class="slide-inner">
                    <span class="slide-label">TIMELINE</span>
                    <h4>10-Week Plan</h4>
                    <div class="divider"></div>
                    <div class="timeline-row">
                        <div class="tl-step" style="--c:#e94560"><span>Wk 1-2<br><strong style="color:#e94560">Discovery</strong></span></div>
                        <div class="tl-step" style="--c:#e94560"><span>Wk 3-4<br><strong style="color:#e94560">Concepts</strong></span></div>
                        <div class="tl-step" style="--c:#e94560"><span>Wk 5-6<br><strong style="color:#e94560">Refine</strong></span></div>
                        <div class="tl-step" style="--c:#e94560"><span>Wk 7-8<br><strong style="color:#e94560">Build</strong></span></div>
                        <div class="tl-step" style="--c:#e94560"><span>Wk 9-10<br><strong style="color:#e94560">Handoff</strong></span></div>
                    </div>
                </div>
                <div class="slide-num">10</div>
            </div>

            <!-- Slide 11: Investment -->
            <div class="slide theme-corp-alt">
                <div class="slide-deco">
                    <div class="corner tl" style="color:#e94560"></div>
                    <div class="corner br" style="color:#e94560"></div>
                    <div class="dots" style="left:16px;bottom:16px;color:#e94560"><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span></div>
                </div>
                <div class="slide-inner">
                    <span class="slide-label">INVESTMENT</span>
                    <h4>₹1,85,000</h4>
                    <div class="divider"></div>
                    <div class="metric-row">
                        <div class="metric"><div class="mv" style="color:#e94560">50%</div><div class="ml">Upfront</div></div>
                        <div class="metric"><div class="mv" style="color:#e94560">25%</div><div class="ml">Concepts</div></div>
                        <div class="metric"><div class="mv" style="color:#e94560">25%</div><div class="ml">Delivery</div></div>
                    </div>
                    <p style="margin-top:8px;font-size:10px">Includes 3 revision rounds per phase</p>
                </div>
                <div class="slide-num">11</div>
            </div>

            <!-- Slide 12: Thank You -->
            <div class="slide theme-corp-dark">
                <div class="slide-deco">
                    <div class="corner tl" style="color:#e94560"></div>
                    <div class="corner br" style="color:#e94560"></div>
                    <div class="circle" style="width:240px;height:240px;left:50%;top:50%;transform:translate(-50%,-50%);background:#e94560;opacity:.04"></div>
                </div>
                <div class="slide-inner slide-end">
                    <h4>Let's Create<br>Something Beautiful</h4>
                    <div class="divider" style="margin:10px auto"></div>
                    <span class="contact">hello@bloomstudio.in<br>+91 98765 43210</span>
                    <span class="tag" style="margin-top:12px">BLOOM STUDIO</span>
                </div>
                <div class="slide-num">12</div>
            </div>
        </div>
    </div>
    <div class="proposal-foot">
        <div class="pf-tags">
            <span class="pf-tag">📄 12 slides</span>
            <span class="pf-tag">🏢 Corporate</span>
            <span class="pf-tag">🔤 Aptos</span>
        </div>
        <a href="/">Recreate this →</a>
    </div>
</div>

<!-- ═══ PROPOSAL 2: Wedding Decor ═══ -->
<div class="proposal">
    <div class="proposal-label" style="color:#f5c6a5">● PROPOSAL 02</div>
    <div class="proposal-head theme-wedding">
        <div class="ph-left">
            <h3>Wedding Decor — Lotus Theme Collection</h3>
            <div class="ph-meta">
                <span>Bloom Studio → Priya &amp; Arjun</span>
                <span class="sep">·</span>
                <span>Jan 22, 2026</span>
                <span class="sep">·</span>
                <span>8 slides</span>
            </div>
        </div>
        <div class="ph-right">
            <div class="price">₹95,000</div>
            <span class="status accepted">Accepted</span>
        </div>
    </div>
    <div class="scroll-hint"><span>Scroll slides</span> <span class="arrow">→</span></div>
    <div class="slides-wrapper">
        <div class="slides-scroll">
            <div class="slide theme-wed">
                <div class="slide-deco">
                    <div class="corner tl" style="color:#f5c6a5"></div>
                    <div class="corner br" style="color:#f5c6a5"></div>
                    <div class="circle" style="width:180px;height:180px;right:-50px;bottom:-50px;background:#c62a88"></div>
                </div>
                <div class="slide-inner slide-title">
                    <span class="tag">CREATIVE PITCH</span>
                    <h4 style="color:#f5c6a5;margin-top:12px">Lotus Theme<br>Collection</h4>
                    <div class="divider" style="background:linear-gradient(90deg,#f5c6a5,#c62a88);margin:8px auto"></div>
                    <p style="color:#eee;font-size:12px">A bespoke wedding experience for <strong>Priya &amp; Arjun</strong></p>
                    <span class="company-from" style="color:#d4a9b8">BY BLOOM STUDIO</span>
                </div>
                <div class="slide-num">01</div>
            </div>
            <div class="slide theme-wed-alt">
                <div class="slide-deco"><div class="corner tl" style="color:#f5c6a5"></div><div class="dots" style="right:16px;bottom:16px;color:#f5c6a5"><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span></div></div>
                <div class="slide-inner">
                    <span class="slide-label" style="color:#f5c6a5">YOUR VISION</span>
                    <h4>Modern Tradition</h4>
                    <div class="divider"></div>
                    <p>A celebration that honors tradition while feeling intimate, modern, and uniquely yours. Lotus symbolism woven throughout every detail.</p>
                </div>
                <div class="slide-num">02</div>
            </div>
            <div class="slide theme-wed-accent">
                <div class="slide-deco"><div class="corner tl" style="color:#f5c6a5"></div><div class="corner br" style="color:#f5c6a5"></div></div>
                <div class="slide-inner">
                    <span class="slide-label" style="color:#f5c6a5">CENTREPIECE</span>
                    <h4>Mandap Design</h4>
                    <div class="divider" style="background:linear-gradient(90deg,#f5c6a5,#c62a88)"></div>
                    <p>Hand-carved wooden structure with cascading lotus garlands, silk draping in ivory and blush, and hanging brass diyas creating warm, ambient glow.</p>
                </div>
                <div class="slide-num">03</div>
            </div>
            <div class="slide theme-wed">
                <div class="slide-deco"><div class="circle" style="width:100px;height:100px;left:-30px;top:-30px;background:#c62a88"></div></div>
                <div class="slide-inner">
                    <span class="slide-label" style="color:#f5c6a5">DETAILS</span>
                    <h4>Table Settings</h4>
                    <div class="divider" style="background:linear-gradient(90deg,#f5c6a5,#c62a88)"></div>
                    <ul class="slide-list">
                        <li style="--c:#f5c6a5">Brass lotus candle holders</li>
                        <li style="--c:#f5c6a5">Hand-painted menu cards</li>
                        <li style="--c:#f5c6a5">Silk Rajasthani block-print runners</li>
                        <li style="--c:#f5c6a5">Fresh flower rangoli</li>
                    </ul>
                </div>
                <div class="slide-num">04</div>
            </div>
            <div class="slide theme-wed-alt">
                <div class="slide-deco"><div class="corner tl" style="color:#f5c6a5"></div></div>
                <div class="slide-inner">
                    <span class="slide-label" style="color:#f5c6a5">ENTRANCE</span>
                    <h4>Welcome Arch</h4>
                    <div class="divider" style="background:linear-gradient(90deg,#f5c6a5,#c62a88)"></div>
                    <p>12-foot arch with woven jasmine and marigold base, oversized paper lotus blooms, and warm LED fairy lights creating a magical first impression.</p>
                </div>
                <div class="slide-num">05</div>
            </div>
            <div class="slide theme-wed-accent">
                <div class="slide-deco"><div class="dots" style="left:16px;bottom:16px;color:#f5c6a5"><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span></div></div>
                <div class="slide-inner">
                    <span class="slide-label" style="color:#f5c6a5">AMBIENCE</span>
                    <h4>Lighting Design</h4>
                    <div class="divider" style="background:linear-gradient(90deg,#f5c6a5,#c62a88)"></div>
                    <p>Warm amber uplighting, floating lotus candles in water features, vintage brass lanterns along pathways, and canopy fairy lights.</p>
                </div>
                <div class="slide-num">06</div>
            </div>
            <div class="slide theme-wed">
                <div class="slide-deco"><div class="corner tl" style="color:#f5c6a5"></div><div class="corner br" style="color:#f5c6a5"></div></div>
                <div class="slide-inner">
                    <span class="slide-label" style="color:#f5c6a5">INVESTMENT</span>
                    <h4>₹95,000</h4>
                    <div class="divider" style="background:linear-gradient(90deg,#f5c6a5,#c62a88)"></div>
                    <p>Complete package including setup, all materials, lighting, day-of coordination, and breakdown. Travel within Jaipur included.</p>
                    <div class="metric-row">
                        <div class="metric"><div class="mv" style="color:#f5c6a5">40%</div><div class="ml">Booking</div></div>
                        <div class="metric"><div class="mv" style="color:#f5c6a5">60%</div><div class="ml">Event Day</div></div>
                    </div>
                </div>
                <div class="slide-num">07</div>
            </div>
            <div class="slide theme-wed-accent">
                <div class="slide-deco">
                    <div class="corner tl" style="color:#f5c6a5"></div><div class="corner br" style="color:#f5c6a5"></div>
                    <div class="circle" style="width:200px;height:200px;left:50%;top:50%;transform:translate(-50%,-50%);background:#c62a88;opacity:.06"></div>
                </div>
                <div class="slide-inner slide-end">
                    <h4 style="color:#f5c6a5">Let's Make Your<br>Day Magical ✨</h4>
                    <div class="divider" style="background:linear-gradient(90deg,#f5c6a5,#c62a88);margin:10px auto"></div>
                    <span class="contact" style="color:#d4a9b8">hello@bloomstudio.in<br>+91 98765 43210</span>
                </div>
                <div class="slide-num">08</div>
            </div>
        </div>
    </div>
    <div class="proposal-foot">
        <div class="pf-tags">
            <span class="pf-tag">📄 8 slides</span>
            <span class="pf-tag">💒 Creative Pitch</span>
            <span class="pf-tag">🌸 Warm Tone</span>
        </div>
        <a href="/">Recreate this →</a>
    </div>
</div>

<!-- ═══ PROPOSAL 3: Corporate Workshop ═══ -->
<div class="proposal">
    <div class="proposal-label" style="color:#64ffda">● PROPOSAL 03</div>
    <div class="proposal-head theme-workshop">
        <div class="ph-left">
            <h3>Corporate Art Workshop — Q1 Team Building</h3>
            <div class="ph-meta">
                <span>Bloom Studio → TechNova Solutions</span>
                <span class="sep">·</span>
                <span>Feb 14, 2026</span>
                <span class="sep">·</span>
                <span>6 slides</span>
            </div>
        </div>
        <div class="ph-right">
            <div class="price">₹45,000</div>
            <span class="status draft">Draft</span>
        </div>
    </div>
    <div class="scroll-hint"><span>Scroll slides</span> <span class="arrow">→</span></div>
    <div class="slides-wrapper">
        <div class="slides-scroll">
            <div class="slide theme-wk">
                <div class="slide-deco">
                    <div class="corner tl" style="color:#64ffda"></div>
                    <div class="corner br" style="color:#64ffda"></div>
                    <div class="circle" style="width:180px;height:180px;right:-50px;bottom:-50px;background:#64ffda"></div>
                </div>
                <div class="slide-inner slide-title">
                    <span class="tag">CORPORATE PROPOSAL</span>
                    <h4 style="color:#64ffda;margin-top:12px">Art &amp; Team<br>Workshop</h4>
                    <div class="divider" style="background:#64ffda;margin:8px auto"></div>
                    <p style="color:#ccd6f6;font-size:12px">Prepared for <strong>TechNova Solutions</strong></p>
                    <span class="company-from" style="color:#64ffda">Q1 TEAM BUILDING</span>
                </div>
                <div class="slide-num">01</div>
            </div>
            <div class="slide theme-wk-alt">
                <div class="slide-deco"><div class="corner tl" style="color:#64ffda"></div><div class="dots" style="right:16px;bottom:16px;color:#64ffda"><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span><span></span></div></div>
                <div class="slide-inner">
                    <span class="slide-label" style="color:#64ffda">WHY ART?</span>
                    <h4>Creative Impact</h4>
                    <div class="divider" style="background:#64ffda"></div>
                    <p>Creative activities boost lateral thinking, reduce stress, and build trust between team members in ways traditional offsites can't match.</p>
                    <div class="metric-row">
                        <div class="metric"><div class="mv" style="color:#64ffda">40%</div><div class="ml">Creativity ↑</div></div>
                        <div class="metric"><div class="mv" style="color:#64ffda">67%</div><div class="ml">Team Trust ↑</div></div>
                    </div>
                </div>
                <div class="slide-num">02</div>
            </div>
            <div class="slide theme-wk-accent">
                <div class="slide-deco"><div class="corner tl" style="color:#64ffda"></div><div class="corner br" style="color:#64ffda"></div></div>
                <div class="slide-inner">
                    <span class="slide-label" style="color:#64ffda">THE EXPERIENCE</span>
                    <h4>3-Hour Session</h4>
                    <div class="divider" style="background:#64ffda"></div>
                    <div class="timeline-row">
                        <div class="tl-step"><span>45 min<br><strong style="color:#64ffda">Basics</strong></span></div>
                        <div class="tl-step"><span>60 min<br><strong style="color:#64ffda">Guided</strong></span></div>
                        <div class="tl-step"><span>75 min<br><strong style="color:#64ffda">Mural</strong></span></div>
                        <div class="tl-step"><span>Gallery<br><strong style="color:#64ffda">Walk</strong></span></div>
                    </div>
                </div>
                <div class="slide-num">03</div>
            </div>
            <div class="slide theme-wk">
                <div class="slide-deco"><div class="circle" style="width:120px;height:120px;left:-30px;bottom:-30px;background:#64ffda"></div></div>
                <div class="slide-inner">
                    <span class="slide-label" style="color:#64ffda">INCLUDED</span>
                    <h4>What's Provided</h4>
                    <div class="divider" style="background:#64ffda"></div>
                    <ul class="slide-list">
                        <li style="--c:#64ffda">All materials (brushes, paints, canvas)</li>
                        <li style="--c:#64ffda">2 professional art instructors</li>
                        <li style="--c:#64ffda">Aprons for all participants</li>
                        <li style="--c:#64ffda">4×6ft collaborative canvas</li>
                        <li style="--c:#64ffda">Framing for individual works</li>
                    </ul>
                </div>
                <div class="slide-num">04</div>
            </div>
            <div class="slide theme-wk-alt">
                <div class="slide-deco"><div class="corner tl" style="color:#64ffda"></div><div class="corner br" style="color:#64ffda"></div></div>
                <div class="slide-inner">
                    <span class="slide-label" style="color:#64ffda">INVESTMENT</span>
                    <h4>₹45,000</h4>
                    <div class="divider" style="background:#64ffda"></div>
                    <p>For up to 25 participants. Includes all materials, instruction, venue setup, and a finished mural for your office wall.</p>
                    <div class="metric-row">
                        <div class="metric"><div class="mv" style="color:#64ffda">25</div><div class="ml">Max People</div></div>
                        <div class="metric"><div class="mv" style="color:#64ffda">3hr</div><div class="ml">Duration</div></div>
                    </div>
                </div>
                <div class="slide-num">05</div>
            </div>
            <div class="slide theme-wk-accent">
                <div class="slide-deco">
                    <div class="corner tl" style="color:#64ffda"></div><div class="corner br" style="color:#64ffda"></div>
                    <div class="circle" style="width:200px;height:200px;left:50%;top:50%;transform:translate(-50%,-50%);background:#64ffda;opacity:.04"></div>
                </div>
                <div class="slide-inner slide-end">
                    <h4 style="color:#64ffda">Let's Create<br>Together 🎨</h4>
                    <div class="divider" style="background:#64ffda;margin:10px auto"></div>
                    <span class="contact" style="color:#8892b0">hello@bloomstudio.in<br>+91 98765 43210</span>
                    <span class="tag" style="margin-top:12px">BLOOM STUDIO</span>
                </div>
                <div class="slide-num">06</div>
            </div>
        </div>
    </div>
    <div class="proposal-foot">
        <div class="pf-tags">
            <span class="pf-tag">📄 6 slides</span>
            <span class="pf-tag">🏢 Corporate</span>
            <span class="pf-tag">🎨 Creative</span>
        </div>
        <a href="/">Recreate this →</a>
    </div>
</div>

</div>

<style>
.slide-list li::before{background:var(--c, #888)}
.tl-step::before{background:var(--c, #888)}
.tl-step::after{background:var(--c, #888)}
</style>
</body></html>'''
